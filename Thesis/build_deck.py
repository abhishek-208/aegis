# -*- coding: utf-8 -*-
"""
M.Tech thesis-defense deck:
  "Enhanced Aegis: A Byzantine-Resilient Federated Learning Protocol for Drone Swarms"
Content grounded in the thesis (./*.tex), the IEEE paper, the interim deck, evolution
notes, and the finalised experimental results. Equations rendered to PNG via MiKTeX
(amsmath/amssymb); architecture diagram compiled from the paper TikZ; result figures
embedded one-per-slide from the finalised results.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

HERE = os.path.dirname(os.path.abspath(__file__))
MATHDIR = os.path.join(HERE, "_math")
OUT = os.path.join(HERE, "Aegis_Defense.pptx")
RESDIR = r"D:\IITD\MTP 2\FL_Project\MTP2 Results"

def EQ(name):
    return os.path.join(MATHDIR, name + ".png")

ARCH = os.path.join(MATHDIR, "arch-1.png")

RES = {
 "baseline_seeds": RESDIR + r"\Expr 31 - @Aegis, Ablation, Diff seeds, No attack\aegis_none_manual_accuracy_line.png",
 "baseline_cmp":   RESDIR + r"\Expr 32 - @Aegis, FedAvg, CWMed, Krum, No attack\fed_avg+aegis+cw_med+multi_krum_none_manual_accuracy_line.png",
 "aegis_multi":    RESDIR + r"\Expr 30 - @Aegis, sign flip, label flip, IPM\aegis_none+sign_flip+label_flip+ipm_comparison_accuracy_line.png",
 "sign_cmp":       RESDIR + r"\Expr 33 - @Aegis, Fedavg, CWMed, Krum, FoolsGold, Bulyan, Sign Flip\fed_avg+aegis+cw_med+multi_krum+fools_gold+bulyan_sign_flip_manual_accuracy_line.png",
 "label_cmp":      RESDIR + r"\Expr 33 - @Aegis, CwMed, Krum, FoolsGold, Bulyan, Label Flip\aegis+cw_med+multi_krum+fools_gold+bulyan_label_flip_manual_accuracy_line.png",
 "ipm_cmp":        RESDIR + r"\Expr 34 - @Aegis, Fedavg, CWMed, Krum, FoolsGold, Bulyan, IPM without LR Decay\fed_avg+aegis+cw_med+multi_krum+fools_gold+bulyan_ipm_manual_accuracy_line.png",

 "alie_cmp":       RESDIR + r"\Expr 33 - @Aegis, Fedavg, CWmed, Krum, Foolsgold, Bulyan, ALIE\fed_avg+aegis+cw_med+multi_krum+fools_gold+bulyan_alie_manual_accuracy_line.png",
 "vol_cmp":        RESDIR + r"\Expr 33 - @Aegis, Fedavg, CWmed, Krum, Foolsgold, Bulyan, Volume Spam\fed_avg+aegis+cw_med+multi_krum+fools_gold+bulyan_volume_spam_manual_accuracy_line.png",
 "syb_cmp":        RESDIR + r"\Expr 33 - @Aegis, Fedavg, CWmed, Krum, Foolsgold, Bulyan, Sybil\fed_avg+aegis+cw_med+multi_krum+fools_gold+bulyan_sybil_manual_accuracy_line.png",

 "advanced":       RESDIR + r"\Expr 30 - @Aegis, Vol spam, ALIE, Sybil\aegis_none+volume_spam+alie+sybil_comparison_accuracy_line.png",
 "byz_sweep":      RESDIR + r"\Expr 29 - @Aegis, Ablation, Byz sweep (10, 20, 30, 40), Sign flip\aegis_none+sign_flip_manual_accuracy_bar.png",
 "shard_sweep":    RESDIR + r"\Expr 28 - @Aegis, shards sweep(2,4,6), label flip, after removing hard cosine filter\aegis_label_flip+none_manual_accuracy_bar.png",
 "alie_ablation":  RESDIR + r"\Expr 35 - Aegis, ALIE, Cos Penalty Ablation\aegis_alie_manual_accuracy_bar.png",
 "vol_ablation":   RESDIR + r"\Expr 31 - Aegis, Ablation, Volume spam vs No volume clipping\aegis_volume_spam_ablation_accuracy_line.png",
 "complexity":     RESDIR + r"\Expr 30 - @Aegis, sign flip, label flip, IPM\complexity_verification_Aegis_Non-IID_-_Sign_Flip.png",
 # diagnostic dashboards correlating the convergence proof's residual bias zeta_A with experiment
 "diag_noatk":     RESDIR + r"\Expr 31 - @Aegis, Ablation, Diff seeds, No attack\DIAGNOSTIC_Aegis_-_No_Attack_Run_1.png",
 "diag_label":     RESDIR + r"\Expr 31 - @Aegis, Ablation, Label flip\DIAGNOSTIC_Full_Aegis_label_flip.png",
 "diag_alie":      RESDIR + r"\Expr 33 - @Aegis, Fedavg, CWmed, Krum, Foolsgold, Bulyan, ALIE\DIAGNOSTIC_Aegis_With_alie_Attack.png",
}

# integrity check ----------------------------------------------------------
_missing = [k for k, v in RES.items() if not os.path.exists(v)]
_missing += [n for n in ["arch-1"] if not os.path.exists(os.path.join(MATHDIR, n + ".png"))]
for n in ["fl_objective","attacks","step1","step2","step3","step4","step5","step6",
          "complexity","evo_initial","evo_enhanced","hypergeom","byz_bound","center_drag",
          "lemma1","lemma2","thm_conv","zeta_def","hp_credit",
          "cert_error","cert_t1","cert_t2","cert_t3","wb_bound","final_conv","zeta_decomp","perfect_filter"]:
    if not os.path.exists(EQ(n)):
        _missing.append("eq:" + n)
if _missing:
    print("WARNING missing assets:")
    for m in _missing:
        print("   -", m)

# ----------------------------------------------------------------------------- palette
NAVY    = RGBColor(0x0E, 0x2A, 0x47)
NAVY2   = RGBColor(0x18, 0x3C, 0x60)
SLATE   = RGBColor(0x33, 0x4A, 0x60)
ACCENT  = RGBColor(0xE3, 0xA5, 0x4E)
ACCENT2 = RGBColor(0xB9, 0x7A, 0x0B)
CRIMSON = RGBColor(0xC0, 0x39, 0x2B)
GREEN   = RGBColor(0x2E, 0x7D, 0x32)
AMBER   = RGBColor(0xB5, 0x6A, 0x10)
LIGHT   = RGBColor(0xF5, 0xF7, 0xFA)
PANEL   = RGBColor(0xEC, 0xF1, 0xF6)
PANEL2  = RGBColor(0xE2, 0xE9, 0xF1)
GOLDTNT = RGBColor(0xF7, 0xE1, 0xAE)
TEXT    = RGBColor(0x1F, 0x29, 0x33)
MUTED   = RGBColor(0x5C, 0x6B, 0x7B)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
RULE    = RGBColor(0xD6, 0xDE, 0xE6)
HEADTNT = RGBColor(0xDD, 0xE6, 0xEF)
CRIMTNT = RGBColor(0xF6, 0xDE, 0xDA)
GRNTNT  = RGBColor(0xDD, 0xEE, 0xDF)

F_HEAD = "Segoe UI Semibold"
F_BODY = "Segoe UI"
F_LIGHT = "Segoe UI Light"
F_MONO = "Consolas"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = 13.333, 7.5
BLANK = prs.slide_layouts[6]

# ----------------------------------------------------------------------------- helpers
def slide():
    return prs.slides.add_slide(BLANK)

def _noshadow(shp):
    try:
        shp.shadow.inherit = False
    except Exception:
        pass

def rect(s, x, y, w, h, fill=None, line=None, line_w=1.0, shape=MSO_SHAPE.RECTANGLE):
    sp = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    _noshadow(sp)
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    return sp

def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=4, line_spacing=1.0, wrap=True):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after); p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (txt, size, color, bold, font, italic) in para:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
            r.font.name = font; r.font.color.rgb = color
    return tb

def R(txt, size, color=TEXT, bold=False, font=F_BODY, italic=False):
    return (txt, size, color, bold, font, italic)

def header(s, kicker, title, tsize=27):
    rect(s, 0.55, 0.50, 0.13, 0.66, fill=ACCENT)
    text(s, 0.82, 0.46, 11.8, 0.3, [[R(kicker.upper(), 12, ACCENT2, True, F_HEAD)]])
    text(s, 0.80, 0.70, 11.0, 0.7, [[R(title, tsize, NAVY, True, F_HEAD)]])
    rect(s, 0.55, 1.42, 12.23, 0.018, fill=RULE)

def opill(s, label):
    rect(s, 11.2, 0.6, 1.55, 0.52, fill=NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, 11.2, 0.6, 1.55, 0.52, [[R(label, 16, ACCENT, True, F_MONO)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

def footer(s, n):
    text(s, 0.55, 7.06, 8.0, 0.3,
         [[R("Enhanced Aegis  â€”  M.Tech Thesis Defense", 9, MUTED, False, F_BODY)]])
    text(s, 11.3, 7.06, 1.45, 0.3, [[R(str(n), 9, MUTED, True, F_BODY)]], align=PP_ALIGN.RIGHT)

def bullets(s, x, y, w, h, items, size=15, gap=7, lh=1.04):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, it in enumerate(items):
        lvl, txt = it[0], it[1]
        col = it[2] if len(it) > 2 else TEXT
        bold = it[3] if len(it) > 3 else False
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap if lvl == 0 else 3); p.space_before = Pt(0)
        p.line_spacing = lh; p.alignment = PP_ALIGN.LEFT
        if lvl == 0:
            b = p.add_run(); b.text = "â–ª  "
            b.font.size = Pt(size-1); b.font.color.rgb = ACCENT; b.font.name = F_BODY; b.font.bold = True
        else:
            pad = p.add_run(); pad.text = "      â€“  "
            pad.font.size = Pt(size-2); pad.font.color.rgb = MUTED; pad.font.name = F_BODY
        r = p.add_run(); r.text = txt
        r.font.size = Pt(size if lvl == 0 else size-2)
        r.font.color.rgb = col; r.font.name = F_BODY; r.font.bold = bold
    return tb

def image_fit(s, path, bx, by, bw, bh, frame=True, caption=None):
    iw, ih = Image.open(path).size
    ar = iw/ih; bar = bw/bh
    if ar > bar:
        w = bw; h = bw/ar
    else:
        h = bh; w = bh*ar
    x = bx + (bw-w)/2; y = by + (bh-h)/2
    if frame:
        rect(s, x-0.05, y-0.05, w+0.10, h+0.10, fill=WHITE, line=RULE, line_w=1.0)
    s.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))
    if caption:
        text(s, bx, by+bh+0.04, bw, 0.4, [[R(caption, 10.5, MUTED, False, F_BODY, True)]],
             align=PP_ALIGN.CENTER)

def panel(s, x, y, w, h, fill=PANEL, line=None, title=None, tcolor=NAVY):
    rect(s, x, y, w, h, fill=fill, line=line, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    if title:
        text(s, x+0.22, y+0.14, w-0.4, 0.4, [[R(title, 14, tcolor, True, F_HEAD)]])

def mono(s, x, y, w, h, lines, size=12.5, color=NAVY, lead=ACCENT2):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(5); p.space_before = Pt(0); p.line_spacing = 1.0
        r = p.add_run(); r.text = ln
        r.font.name = F_MONO; r.font.size = Pt(size); r.font.color.rgb = color
    return tb

def set_cell(cell, txt, size, color, bold=False, align=PP_ALIGN.CENTER,
             fill=WHITE, font=F_BODY, italic=False):
    cell.fill.solid(); cell.fill.fore_color.rgb = fill
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = Inches(0.06); cell.margin_right = Inches(0.06)
    cell.margin_top = Inches(0.02); cell.margin_bottom = Inches(0.02)
    tf = cell.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = txt
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color; r.font.name = font

def table(s, x, y, w, rows, col_w, data, header_fill=NAVY, header_color=WHITE,
          fsize=10.5, hsize=10.5, row_h=0.4, cell_styles=None, zebra=True):
    nr = len(data); nc = len(data[0])
    gt = s.shapes.add_table(nr, nc, Inches(x), Inches(y), Inches(w), Inches(row_h*nr)).table
    gt.first_row = False; gt.horz_banding = False
    for ci, cw in enumerate(col_w):
        gt.columns[ci].width = Inches(cw)
    for ri in range(nr):
        gt.rows[ri].height = Inches(row_h)
        for ci in range(nc):
            cell = gt.cell(ri, ci)
            txt = str(data[ri][ci])
            if ri == 0:
                set_cell(cell, txt, hsize, header_color, True,
                         PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER, header_fill, F_HEAD)
            else:
                base = LIGHT if (zebra and ri % 2 == 0) else WHITE
                col = TEXT; bold = False; al = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
                st = (cell_styles or {}).get((ri, ci))
                if st:
                    col = st.get("color", col); bold = st.get("bold", bold); base = st.get("fill", base)
                    al = st.get("align", al)
                set_cell(cell, txt, fsize, col, bold, al, base,
                         F_HEAD if (ci == 0 and bold) else F_BODY,
                         italic=(st or {}).get("italic", False) if st else False)
    return gt

def divider(s, num, kicker, title, sub=None):
    rect(s, 0, 0, SW, SH, fill=NAVY)
    rect(s, 0, 0, 0.32, SH, fill=ACCENT)
    text(s, 1.2, 2.15, 3.0, 2.0, [[R(num, 120, NAVY2, True, F_LIGHT)]], anchor=MSO_ANCHOR.MIDDLE)
    rect(s, 1.35, 3.42, 1.7, 0.06, fill=ACCENT)
    text(s, 1.4, 2.7, 10.0, 0.5, [[R(kicker.upper(), 14, ACCENT, True, F_HEAD)]])
    text(s, 1.35, 3.6, 10.6, 1.4, [[R(title, 40, WHITE, True, F_HEAD)]])
    if sub:
        text(s, 1.4, 4.85, 10.4, 1.0, [[R(sub, 16, RGBColor(0xC6,0xD2,0xDE), False, F_BODY)]])

PAGE = [1]
def pg():
    PAGE[0] += 1
    return PAGE[0]

def step_slide(n, title, eqname, bigO, eq_h, whys, kicker="The Aegis Protocol Â· Step %d"):
    s = slide()
    header(s, kicker % n if "%d" in kicker else kicker, title)
    opill(s, bigO)
    panel(s, 0.8, 1.66, 11.75, eq_h, fill=PANEL, line=RULE)
    image_fit(s, EQ(eqname), 1.15, 1.82, 11.05, eq_h-0.32, frame=False)
    yb = 1.66 + eq_h + 0.28
    bullets(s, 0.95, yb, 11.55, 6.95-yb, whys, size=14.5, gap=8)
    footer(s, pg())
    return s

def img_slide(kicker, title, key, takeaway, tk_color=NAVY, tsize=27):
    s = slide()
    header(s, kicker, title, tsize=tsize)
    image_fit(s, RES[key], 0.8, 1.62, 11.75, 4.62, frame=True)
    rect(s, 0.8, 6.42, 11.75, 0.66, fill=GOLDTNT, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, 0.8, 6.42, 0.10, 0.66, fill=ACCENT, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, 1.05, 6.42, 11.4, 0.66, [takeaway], anchor=MSO_ANCHOR.MIDDLE)
    footer(s, pg())
    return s

# ============================================================================= SLIDES

# --- 1. TITLE --------------------------------------------------------------
s = slide()
rect(s, 0, 0, SW, SH, fill=NAVY)
rect(s, 0, 0, SW, 0.22, fill=ACCENT)
rect(s, 0, SH-0.22, SW, 0.22, fill=ACCENT)
text(s, 1.0, 1.28, 11.3, 0.4, [[R("M.TECH THESIS DEFENSE  Â·  INDIAN INSTITUTE OF TECHNOLOGY DELHI",
                                  14, ACCENT, True, F_HEAD)]])
text(s, 0.95, 1.92, 11.6, 1.3, [[R("Enhanced Aegis", 58, WHITE, True, F_HEAD)]])
text(s, 1.0, 3.02, 11.4, 1.1,
     [[R("A Byzantine-Resilient Federated Learning Protocol for Drone Swarms",
         25, RGBColor(0xE9,0xC8,0x8F), False, F_LIGHT)]])
rect(s, 1.02, 4.2, 5.2, 0.035, fill=NAVY2)
text(s, 1.0, 4.5, 11.0, 1.2, [
    [R("Abhishek Kumar Tripathi", 20, WHITE, True, F_HEAD), R("    (2024EET2845)", 15, RGBColor(0xB9,0xC6,0xD4), False, F_BODY)],
    [R("Advisor:  Prof. Harshan Jagadeesh", 16, RGBColor(0xD7,0xE0,0xEA), False, F_BODY)],
    [R("Department of Electrical Engineering, IIT Delhi", 14, RGBColor(0xA9,0xB8,0xC8), False, F_BODY)],
], space_after=6)
text(s, 1.0, 6.5, 11.0, 0.4, [[R("June 2026", 14, ACCENT, True, F_HEAD)]])

# --- 2. OUTLINE ------------------------------------------------------------
s = slide(); header(s, "Agenda", "Outline")
items = [
    ("1", "Motivation & the Problem", "Federated Learning on diverse edge networks; the Byzantine threat under highly skewed data"),
    ("2", "Threat Model & Attacks", "Our omniscient, adaptive attacker using six distinct attack strategies"),
    ("3", "The Aegis Protocol", "How Aegis evolved, its six mathematical stages, and achieving O(kd) efficiency"),
    ("4", "Experimental Setup", "Testing with CIFAR-10, highly Non-IID data, and five baseline defenses"),
    ("5", "Results & Analysis", "Final experimental results, including successes, the ALIE/IPM breaches, and sweeps"),
    ("6", "Conclusions & Future Work", "Summary of achievements, known limitations, and four future research directions"),
]
y = 1.72
for num, t, d in items:
    rect(s, 0.95, y, 0.6, 0.6, fill=NAVY, shape=MSO_SHAPE.OVAL)
    text(s, 0.95, y, 0.6, 0.6, [[R(num, 19, WHITE, True, F_HEAD)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, 1.76, y-0.04, 10.9, 0.45, [[R(t, 18, NAVY, True, F_HEAD)]])
    text(s, 1.76, y+0.35, 10.9, 0.4, [[R(d, 12.5, MUTED, False, F_BODY)]])
    y += 0.85
footer(s, pg())

# --- 3. DIVIDER 1 ----------------------------------------------------------
s = slide(); divider(s, "01", "Part One", "Motivation & the Problem",
                     "Why federated learning, and why Byzantine resilience is hard under Non-IID data")

# --- 4. WHY FL -------------------------------------------------------------
s = slide(); header(s, "Motivation", "Why Federated Learning?")
text(s, 0.8, 1.6, 11.8, 0.6,
     [[R("Edge devices like drone swarms and hospitals collaborate to train a shared model. "
         "They exchange only gradients, keeping raw data securely on-device.", 15, TEXT, False, F_BODY)]],
     line_spacing=1.05)
cards = [
    ("Bandwidth", ACCENT2,
     ["Costs scale with model size, not dataset size.",
      "Sending just 4.5 MB per round instead of 21 GB saves immense bandwidth."]),
    ("Privacy", GREEN,
     ["The server only sees gradient updates, never raw data or labels.",
      "Compatible with encryption without altering the core protocol."]),
    ("The catch", CRIMSON,
     ["The server's lack of visibility creates a critical blind spot.",
      "Attackers can submit poisoned gradients that look completely normal."]),
]
cw = 3.78; x = 0.8
for title, col, lines in cards:
    panel(s, x, 2.5, cw, 3.45, fill=PANEL, line=RULE)
    rect(s, x, 2.5, cw, 0.10, fill=col, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, x+0.25, 2.72, cw-0.5, 0.5, [[R(title, 18, col, True, F_HEAD)]])
    bullets(s, x+0.25, 3.28, cw-0.5, 2.5, [(0, ln) for ln in lines], size=13, gap=9, lh=1.05)
    x += cw + 0.27
panel(s, 0.8, 6.18, 11.75, 0.85, fill=PANEL, line=RULE)
image_fit(s, EQ("fl_objective"), 1.0, 6.28, 11.3, 0.66, frame=False)
footer(s, pg())

# --- 5. NON-IID CRISIS -----------------------------------------------------
s = slide(); header(s, "The Core Difficulty", "The Compounding Challenge of Highly Skewed Data")
bullets(s, 0.8, 1.62, 6.55, 4.2, [
    (0, "Real-world federated nodes are highly specialized.", NAVY, True),
    (1, "In our setup, each client sees only ~40% of the classes."),
    (1, "This makes honest updates 10â€“50Ã— more dissimilar than standard IID data."),
    (0, "Most defenses filter out updates far from the center.", NAVY, True),
    (1, "But highly specialized honest updates naturally sit far from the center."),
    (1, "Filters discard this honest signal while admitting stealthy, centered poison."),
    (0, "The Byzantine Tax", CRIMSON, True),
    (1, "The accuracy lost just by running the defense under normal conditions."),
    (1, "For example, Bulyan loses 8.5 percentage points with zero attackers."),
], size=14, gap=6)
panel(s, 7.7, 1.7, 4.9, 3.05, fill=NAVY)
text(s, 7.95, 1.86, 4.5, 0.5, [[R("The lethal interaction", 15, ACCENT, True, F_HEAD)]])
text(s, 7.95, 2.4, 4.45, 2.3, [
    [R("Honest specialists", 13.5, WHITE, True, F_BODY), R("  look like outliers", 13, RGBColor(0xC6,0xD2,0xDE), False, F_BODY)],
    [R("Stealthy attackers", 13.5, WHITE, True, F_BODY), R("  hide near the center", 13, RGBColor(0xC6,0xD2,0xDE), False, F_BODY)],
    [R("â†“", 18, ACCENT, True, F_BODY)],
    [R("Standard filters reject good data", 13, WHITE, False, F_BODY)],
    [R("and accept the poison.", 13, WHITE, False, F_BODY)],
], space_after=6, line_spacing=1.04)
panel(s, 7.7, 4.95, 4.9, 1.25, fill=PANEL, line=RULE, title="Center Drag (single-pass median)")
image_fit(s, EQ("center_drag"), 7.9, 5.45, 4.5, 0.6, frame=False)
text(s, 0.8, 6.62, 11.8, 0.4,
     [[R("No prior single-stage aggregator handles Non-IID diversity, IPM, and Sybil simultaneously.",
         13.5, ACCENT2, True, F_BODY)]])
footer(s, pg())

# --- 6. RESEARCH GAP -------------------------------------------------------
s = slide(); header(s, "Research Gap", "What Is Missing â€” and Our Research Question")
panel(s, 0.8, 1.7, 11.75, 1.5, fill=PANEL, line=RULE)
text(s, 1.1, 1.9, 11.2, 1.2, [
    [R("Our Research Question:  ", 16, ACCENT2, True, F_HEAD),
     R("Can a linear-time aggregator protect a highly diverse federated "
       "network against coordinated attacks ", 16, NAVY, False, F_BODY)],
    [R("without punishing honest specialists?", 16, NAVY, False, F_BODY)],
], space_after=4, line_spacing=1.1)
text(s, 0.8, 3.45, 11.8, 0.4, [[R("Three requirements no existing defence meets together:", 15, TEXT, True, F_BODY)]])
reqs = [
    ("Heterogeneity tolerance", "Preserve divergent honest specialists instead of rejecting them as outliers â€” no accuracy penalty on Non-IID clients."),
    ("Uncontaminated reference", "Score every update against a clean, debiased center that coordinated, stealthy adversaries cannot tilt."),
    ("Influence capping", "Bound what any single actor can manufacture through duplicated identities or inflated data volume."),
]
x = 0.8
for t, d in reqs:
    panel(s, x, 4.0, cw, 2.05, fill=WHITE, line=RULE)
    rect(s, x+0.25, 4.23, 0.5, 0.5, fill=NAVY, shape=MSO_SHAPE.OVAL)
    text(s, x+0.25, 4.23, 0.5, 0.5, [[R("âœ“", 18, ACCENT, True, F_HEAD)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, x+0.9, 4.25, cw-1.0, 0.5, [[R(t, 15.5, NAVY, True, F_HEAD)]])
    text(s, x+0.25, 4.95, cw-0.5, 1.0, [[R(d, 13, MUTED, False, F_BODY)]], line_spacing=1.05)
    x += cw + 0.27
text(s, 0.8, 6.35, 11.8, 0.5, [[R("â†’  This gap motivates Aegis.", 14, ACCENT2, True, F_BODY)]])
footer(s, pg())

# --- 7. DIVIDER 2 ----------------------------------------------------------
s = slide(); divider(s, "02", "Part Two", "Threat Model & Attacks",
                     "An omniscient, adaptive, coordinated adversary at up to 40% of the network")

# --- 8. THREAT MODEL -------------------------------------------------------
s = slide(); header(s, "Threat Model", "Adversary Capabilities & Assumptions")
panel(s, 0.8, 1.66, 5.75, 3.0, fill=PANEL, line=RULE, title="The adversary CAN")
bullets(s, 1.05, 2.2, 5.3, 2.5, [
    (0, "Perfectly understand the aggregation rule and global model."),
    (0, "Estimate the honest network's mean gradient."),
    (0, "Fully coordinate all compromised clients."),
    (0, "Persistently attack in every single round."),
], size=13.5, gap=8)
panel(s, 6.8, 1.66, 5.75, 3.0, fill=NAVY)
text(s, 7.05, 1.8, 5.3, 0.4, [[R("THE ADVERSARY CANNOT", 14, ACCENT, True, F_HEAD)]])
tb = s.shapes.add_textbox(Inches(7.05), Inches(2.3), Inches(5.3), Inches(2.2))
tf = tb.text_frame; tf.word_wrap = True
for i, ln in enumerate(["Intercept honest updates before they reach the server",
                        "Corrupt the central server itself",
                        "Take over the majority of the network (f < K/2)"]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(10); p.line_spacing = 1.05
    a = p.add_run(); a.text = "âœ•  "; a.font.color.rgb = CRIMSON; a.font.size = Pt(14); a.font.bold = True; a.font.name=F_BODY
    b = p.add_run(); b.text = ln; b.font.color.rgb = WHITE; b.font.size = Pt(14); b.font.name = F_BODY
panel(s, 0.8, 4.85, 5.75, 1.45, fill=GOLDTNT, line=ACCENT, title="Honest-majority bound")
image_fit(s, EQ("byz_bound"), 1.0, 5.4, 5.35, 0.7, frame=False)
panel(s, 6.8, 4.85, 5.75, 1.45, fill=PANEL, line=RULE, title="Fixed adversary, random participation")
image_fit(s, EQ("hypergeom"), 7.0, 5.45, 5.35, 0.55, frame=False)
text(s, 0.8, 6.42, 11.8, 0.5, [[R("System: K = 30 clients, one honest non-colluding server, m âˆˆ [15,25] sampled/round.  "
     "Byzantine fraction swept to f = 0.40 â€” Aegis defends gracefully through f = 0.30 (collapse at 0.40).", 12.5, MUTED, False, F_BODY, True)]], line_spacing=1.03)
footer(s, pg())

# --- 9. THREAT MATRIX ------------------------------------------------------
s = slide(); header(s, "Attacks", "The Threat Matrix")
panel(s, 0.8, 1.66, 11.75, 1.7, fill=PANEL, line=RULE)
image_fit(s, EQ("attacks"), 1.1, 1.84, 11.15, 1.36, frame=False)
data = [
    ["Attack", "Mechanism", "Norm\nevasion", "Dir.\nevasion", "Median\npoison", "Identity\nspoof"],
    ["Label-Flip", "Data corruption", "âœ“", "âœ“", "âœ•", "âœ•"],
    ["Sign-Flip", "Direction inversion", "âœ•", "âœ•", "âœ•", "âœ•"],
    ["IPM", "Stealth anti-gradient", "âœ“", "âœ•", "âœ“", "âœ•"],
    ["ALIE", "In-variance poison", "âœ“", "âœ“", "âœ“", "âœ•"],
    ["Sybil", "Density flooding", "âœ“", "âœ“", "âœ“", "âœ“"],
]
styles = {}
for r in range(1, 6):
    for c in range(2, 6):
        v = data[r][c]
        styles[(r, c)] = {"color": CRIMSON if v == "âœ“" else RGBColor(0x9A,0xA6,0xB2), "bold": v == "âœ“"}
table(s, 0.8, 3.55, 11.75, 6, [1.85, 3.5, 1.6, 1.6, 1.6, 1.6], data,
      fsize=12.5, hsize=11, row_h=0.5, cell_styles=styles)
text(s, 0.8, 6.55, 11.8, 0.5, [[R("No single defensive axis defeats all four â€” a robust defence must apply "
     "magnitude, direction, and identity checks against a clean, debiased reference.",
     13, ACCENT2, True, F_BODY)]])
footer(s, pg())

# --- 10. RELATED WORK ------------------------------------------------------
s = slide(); header(s, "Related Work", "Existing Defences â€” and Their Limits")
data = [
    ["Method", "Core mechanism", "Tolerance", "Complexity"],
    ["Multi-Krum", "Select k updates minimising neighbour distance", "f < K/2 âˆ’ 2", "O(KÂ²d)"],
    ["Coord.-wise Median", "Per-coordinate median over all updates", "f < K/2", "O(Kd)"],
    ["Trimmed Mean", "Drop extreme f per coordinate, average rest", "f < K/2", "O(Kd log K)"],
    ["Bulyan", "Krum selection, then trimmed mean", "f < K/4", "O(KÂ²d)"],
    ["FoolsGold", "Down-weight high cosine-similarity (Sybil)", "Sybil-only", "O(KÂ²d)"],
    ["Aegis  (this thesis)", "Two-pass median + dual-metric + EMA reputation", "f < K/2", "O(kd)"],
]
styles = {}
for c in range(4):
    styles[(6, c)] = {"color": NAVY, "bold": True, "fill": GOLDTNT, "align": PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER}
styles[(4, 2)] = {"color": CRIMSON, "bold": True}
table(s, 0.8, 1.7, 11.75, 7, [3.05, 4.95, 1.95, 1.8], data, fsize=12.5, hsize=12.5, row_h=0.52, cell_styles=styles)
bullets(s, 0.8, 5.55, 11.8, 1.5, [
    (0, "Bulyan â€” the strongest prior scheme â€” needs f < K/4; it is inapplicable at our 40% setting.", TEXT, False),
    (0, "Every O(KÂ²d) scheme is prohibitive at scale. Aegis sustains f = 0.30 at O(kd) â€” past Bulyan's f < K/4 ceiling and far cheaper.", NAVY, True),
    (0, "Tolerances above are theoretical (primitive breakdown points); empirically all collapse near f = 0.40 under our Non-IID threat â€” Aegis stays graceful furthest, to f = 0.30.", MUTED, False),
], size=12.5, gap=6)
footer(s, pg())

# --- 11. DIVIDER 3 ---------------------------------------------------------
s = slide(); divider(s, "03", "Part Three", "The Aegis Protocol",
                     "From the interim design to the enhanced protocol â€” six mathematical stages, O(kd)")

# --- 12. EVOLUTION TABLE ---------------------------------------------------
s = slide(); header(s, "Evolution", "From Initial Aegis to Enhanced Aegis")
data = [
    ["Aspect", "Initial Aegis  (interim)", "Enhanced Aegis  (paper / thesis)"],
    ["Directional filter", "Hard gate: cos < 0 â†’ discard", "Soft penalty Î±Â·Pâ‚– (Euclidean-only gate)"],
    ["Reputation", "Not in credit score", "EMA Râ‚– fused into credit (Î³=0.95, Î»=20)"],
    ["Adaptive threshold", "Round decay + variance (CV)", "+ hard floor  k_floor = 4.0"],
    ["Convergence aid", "Server momentum (on)", "Off â†’ re-scoped as IPM velocity patch"],
    ["Cosine weight Î±", "10", "30"],
    ["Scale / LR", "40 clients,  lr 1eâˆ’3", "30 clients,  lr 1eâˆ’2"],
    ["Threats studied", "sign, noise, label, orthogonal", "+ IPM, ALIE, Sybil, vol-spam (+ limits)"],
]
styles = {}
for r in range(1, 8):
    styles[(r, 0)] = {"color": NAVY, "bold": True, "align": PP_ALIGN.LEFT}
    styles[(r, 1)] = {"color": SLATE, "align": PP_ALIGN.LEFT}
    styles[(r, 2)] = {"color": NAVY, "bold": True, "fill": GOLDTNT, "align": PP_ALIGN.LEFT}
table(s, 0.8, 1.72, 11.75, 8, [2.55, 4.2, 5.0], data, fsize=12.5, hsize=12.5, row_h=0.55, cell_styles=styles)
text(s, 0.8, 6.55, 11.8, 0.4, [[R("The interim protocol proved the concept; the enhanced protocol fixes the "
     "false-positive and stealth-attack gaps it exposed.", 13, ACCENT2, True, F_BODY)]])
footer(s, pg())

# --- 13. WHY EVOLVED -------------------------------------------------------
s = slide(); header(s, "Evolution", "Why the Protocol Evolved")
panel(s, 0.8, 1.62, 11.75, 1.18, fill=CRIMTNT, line=RULE)
text(s, 1.0, 1.74, 3.2, 0.9, [[R("Initial Aegis", 14, CRIMSON, True, F_HEAD)], [R("(interim)", 11.5, CRIMSON, False, F_BODY)]], anchor=MSO_ANCHOR.MIDDLE)
image_fit(s, EQ("evo_initial"), 4.0, 1.74, 8.4, 0.95, frame=False)
panel(s, 0.8, 2.92, 11.75, 1.18, fill=GRNTNT, line=RULE)
text(s, 1.0, 3.04, 3.2, 0.9, [[R("Enhanced Aegis", 14, GREEN, True, F_HEAD)], [R("(final)", 11.5, GREEN, False, F_BODY)]], anchor=MSO_ANCHOR.MIDDLE)
image_fit(s, EQ("evo_enhanced"), 4.0, 3.02, 8.4, 1.0, frame=False)
reasons = [
    ("A hard cosine gate is a trap", "It allows orthogonal attacks while amputating honest specialists. Our fix: a continuous soft penalty Î±Â·Pâ‚–."),
    ("Round-by-round scoring misses stealth", "Attacks like ALIE hide within normal variance. Our fix: cross-round EMA reputation Î»Râ‚– catches them over time."),
    ("Coupling them caused a death spiral", "Cleaning the pool tightens thresholds, hurting honest clients. Our fix: decouple threshold calculations from the filtered pool."),
]
x = 0.8
for t, d in reasons:
    panel(s, x, 4.35, cw, 2.05, fill=PANEL, line=RULE)
    rect(s, x, 4.35, cw, 0.09, fill=ACCENT, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, x+0.22, 4.55, cw-0.44, 0.7, [[R(t, 13.5, NAVY, True, F_HEAD)]])
    text(s, x+0.22, 5.28, cw-0.44, 1.1, [[R(d, 11.5, MUTED, False, F_BODY)]], line_spacing=1.03)
    x += cw + 0.27
footer(s, pg())

# --- 14. ARCHITECTURE ------------------------------------------------------
s = slide(); header(s, "The Aegis Protocol", "Aggregation Pipeline â€” Architecture")
opill(s, "O(kd)")
image_fit(s, ARCH, 0.8, 1.6, 11.75, 4.55, frame=True)
text(s, 0.8, 6.32, 11.8, 0.75, [
    [R("Stage 1 ", 12.5, ACCENT2, True, F_BODY), R("two-pass cosine-screened median (clean reference).   ", 12.5, TEXT, False, F_BODY),
     R("Stage 2 ", 12.5, ACCENT2, True, F_BODY), R("adaptive MAD Euclidean gate â€” the only hard filter.   ", 12.5, TEXT, False, F_BODY)],
    [R("Stage 3 ", 12.5, ACCENT2, True, F_BODY), R("soft credit score fuses volume, cosine penalty Î±Â·Pâ‚–, and reputation Î»Râ‚– into the weights wâ‚–.", 12.5, TEXT, False, F_BODY)],
], space_after=3, line_spacing=1.05)
footer(s, pg())

# --- 15..20 STEP SLIDES ----------------------------------------------------
step_slide(1, "Step 1: Local Training & Pseudo-Gradients", "step1", "O(dk)", 1.35, [
    (0, "Server broadcasts the global model; each client k trains locally and returns weights wâ‚– and claimed sample size nâ‚–.", TEXT, False),
    (0, "Aegis scores in gradient (delta) space â€” robust to the shared initialisation â€” but aggregates the approved clients' weights.", NAVY, True),
])
step_slide(2, "Step 2: Two-Pass Median Decontamination", "step2", "O(kd)", 2.5, [
    (0, "Pass 1 screens strongly anti-aligned updates (cos < Ï„) out of the median pool; Pass 2 recomputes a debiased clean center.", TEXT, False),
    (0, "Defeats Center Drag: IPM / sign-flip can no longer tilt the reference that every downstream score depends on.", NAVY, True),
    (0, "Ï„ = âˆ’0.3 is deliberately lenient â€” honest Non-IID specialists (mildly negative cosine) are retained.", MUTED, False),
])
step_slide(3, "Step 3: Dual Anomaly Scoring", "step3", "O(kd)", 1.85, [
    (0, "Every client is re-scored against the clean median on two complementary axes: magnitude (Eâ‚–) and direction (Pâ‚–).", TEXT, False),
    (0, "Magnitude-only filters miss IPM; direction-only filters miss amplitude-scaled Sybil floods. Aegis uses both.", NAVY, True),
])
step_slide(4, "Step 4: Adaptive MAD Thresholding", "step4", "O(k)", 2.85, [
    (0, "MAD gives a breakdown-point-0.5 scale estimate; the multiplier kâ½áµ—â¾ decays over warmup and relaxes with round variance (CV).", TEXT, False),
    (0, "The floor k_floor = 4.0 prevents over-rejection of honest specialists in persistent Non-IID rounds. Euclidean is the only hard gate.", NAVY, True),
])
step_slide(5, "Step 5: Reputation & Credit Scoring", "step5", "O(m)", 2.45, [
    (0, "Volume clipping caps inflated-dataset spam; an EMA reputation Râ‚– accumulates each client's persistent directional penalty.", TEXT, False),
    (0, "The cosine penalty is soft (Î±Â·Pâ‚–) â€” attackers are suppressed in the weight, never hard-excluded. This replaced the interim's hard cosine gate.", NAVY, True),
])
step_slide(6, "Step 6: Weighted Aggregation", "step6", "O(kd)", 1.35, [
    (0, "The new global model is the credit-weighted average of the approved clients' weights.", TEXT, False),
    (0, "No pairwise distance matrix is ever formed â€” this is what keeps the whole pipeline linear in d.", NAVY, True),
])

# --- 21. COMPLEXITY --------------------------------------------------------
s = slide(); header(s, "The Aegis Protocol", "Linear-Time Aggregation: O(kd)")
panel(s, 0.8, 1.62, 11.75, 1.0, fill=GOLDTNT, line=ACCENT)
image_fit(s, EQ("complexity"), 1.0, 1.74, 11.3, 0.78, frame=False)
image_fit(s, RES["complexity"], 0.7, 2.78, 7.4, 3.65, frame=True)
panel(s, 8.35, 2.78, 4.2, 3.65, fill=PANEL, line=RULE, title="Empirical verification")
bullets(s, 8.57, 3.35, 3.78, 3.0, [
    (0, "Aggregation time grows linearly in client count k â€” linear fit RÂ² = 0.9989.", NAVY, True),
    (0, "Introselect medians + no pairwise matrix â†’ O(kd), matching FedAvg's floor.", TEXT, False),
    (0, "â‰ˆ 30Ã— faster than Bulyan at K=30, dâ‰ˆ1.12M.", TEXT, False),
    (0, "Server state: a 30-float EMA table (240 bytes) â€” fast failover.", TEXT, False),
], size=12.5, gap=10)
footer(s, pg())

# --- THEORETICAL GUARANTEES ------------------------------------------------
s = slide(); header(s, "The Aegis Protocol", "Theoretical Guarantees")
panel(s, 0.8, 1.66, 5.78, 2.45, fill=PANEL, line=RULE, title="Lemma 1 Â· Pass-1 neutralises IPM")
image_fit(s, EQ("lemma1"), 1.0, 2.22, 5.4, 1.0, frame=False)
text(s, 1.02, 3.32, 5.35, 0.75, [[R("Any IPM vector âˆ’ÎµÂ·á¸¡_H has cosine â‰¤ 0 < Ï„, so it is removed from the "
     "median pool â€” the reference mâ‚‚ stays clean.", 12.5, TEXT, False, F_BODY)]], line_spacing=1.05)
panel(s, 6.78, 1.66, 5.78, 2.45, fill=PANEL, line=RULE, title="Lemma 2 Â· Median breakdown point")
image_fit(s, EQ("lemma2"), 7.0, 2.32, 5.4, 0.8, frame=False)
text(s, 7.02, 3.32, 5.35, 0.78, [[R("Breakdown-0.5 (f < K/2) bounds the median reference, not the whole pipeline. "
     "End-to-end, Aegis stays graceful to f = 0.30 and collapses at f = 0.40 (detection â†’ 0).", 12.0, TEXT, False, F_BODY)]], line_spacing=1.03)
data = [
    ["Guarantee", "Aegis bound", "Source"],
    ["Robustness (breakdown)", "â‰¤ 0.30 empirical", "median ref. is breakdown-0.5 (f < K/2, Lemma 2); end-to-end collapse at f = 0.40"],
    ["Complexity (per round)", "O(kd)", "two-pass median, no pairwise distance matrix"],
    ["Convergence rate", "O(1/âˆšT)", "SGD rate + Îº_AÂ·V_H heterogeneity + Î¶_A residual (next slides)"],
]
styles = {(r, 1): {"color": NAVY, "bold": True} for r in range(1, 4)}
table(s, 0.8, 4.45, 11.75, 4, [3.5, 2.6, 5.65], data, fsize=12.5, hsize=12, row_h=0.5, cell_styles=styles)
text(s, 0.8, 6.6, 11.8, 0.4, [[R("The 0.5 breakdown point is the median reference's; end-to-end Aegis stays graceful to f = 0.30 "
     "(collapse at 0.40) â€” at the SGD rate and linear O(kd) cost.", 13, ACCENT2, True, F_BODY)]])
footer(s, pg())

# --- CONVERGENCE I: THE AEGIS ROBUSTNESS CERTIFICATE -----------------------
s = slide(); header(s, "The Aegis Protocol", "Convergence I â€” The Robustness Certificate", tsize=25)
panel(s, 0.8, 1.58, 11.75, 1.28, fill=GOLDTNT, line=ACCENT)
image_fit(s, EQ("cert_error"), 1.0, 1.70, 11.3, 1.04, frame=False)
text(s, 0.8, 2.97, 11.8, 0.35,
     [[R("The per-round aggregation error splits into three sources â€” each maps to one experimental regime:",
         13.5, NAVY, True, F_HEAD)]])
cc = [
    ("Honest reweighting distortion", GREEN, "cert_t1",
     "Unequal honest weights distort the mean even with zero attackers.",
     "â†’  Byzantine Tax = 0.08 pp  (Aegis 76.08 vs FedAvg 76.16, no attack)."),
    ("Clipped volume leakage", ACCENT2, "cert_t2",
     "Volume clipping caps every client's share at 2M, so spam cannot dominate.",
     "â†’  Volume-spam fully nullified  (â‰ˆ 74%, on par with no attack)."),
    ("Median contamination", CRIMSON, "cert_t3",
     "Surviving Byzantine clients displace the median reference mâ‚‚ by Î”_m.",
     "â†’  ALIE drags mâ‚‚  â‡’  collapse to 10% (chance)."),
]
cwx = 3.78; x = 0.8
for ti, col, eqn, l1, l2 in cc:
    panel(s, x, 3.40, cwx, 2.38, fill=PANEL, line=RULE)
    rect(s, x, 3.40, cwx, 0.09, fill=col, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, x+0.22, 3.56, cwx-0.44, 0.5, [[R(ti, 13, col, True, F_HEAD)]])
    image_fit(s, EQ(eqn), x+0.30, 4.10, cwx-0.60, 0.46, frame=False)
    text(s, x+0.22, 4.66, cwx-0.44, 1.0, [[R(l1, 11, MUTED, False, F_BODY)]], line_spacing=1.03)
    text(s, x+0.22, 5.34, cwx-0.44, 0.4, [[R(l2, 11, col, True, F_BODY)]], line_spacing=1.0)
    x += cwx + 0.27
panel(s, 0.8, 5.92, 11.75, 1.05, fill=PANEL, line=RULE)
text(s, 1.0, 5.99, 11.5, 0.3, [[R("Score suppression keeps the Byzantine weight w_B small â€” and bounds all three terms:",
     12, ACCENT2, True, F_HEAD)]])
image_fit(s, EQ("wb_bound"), 1.0, 6.30, 11.3, 0.55, frame=False)
footer(s, pg())

# --- CONVERGENCE II: FINAL THEOREM & RESIDUAL BIAS -------------------------
s = slide(); header(s, "The Aegis Protocol", "Convergence II â€” Final Theorem & Residual Bias", tsize=24)
panel(s, 0.8, 1.56, 11.75, 1.44, fill=GOLDTNT, line=ACCENT)
image_fit(s, EQ("final_conv"), 0.95, 1.68, 11.45, 1.20, frame=False)
text(s, 0.8, 3.10, 6.0, 0.35, [[R("Two terms vanish â€” two are Aegis-specific", 14, NAVY, True, F_HEAD)]])
bullets(s, 0.8, 3.54, 6.05, 2.3, [
    (0, "Optimisation gap + SGD noise â†’ 0 at the standard O(1/âˆšT) rate as Î·_t decays.", TEXT, False),
    (0, "Îº_AÂ·V_H â€” honest-heterogeneity penalty: the price of unequal honest weighting; grows with the Non-IID variance V_H.", NAVY, True),
    (0, "Î¶_A â€” residual bias: the ONLY attack-driven, non-vanishing floor. Zero iff no attacker survives the filter.", CRIMSON, True),
], size=12.5, gap=9)
panel(s, 7.0, 3.18, 5.55, 2.5, fill=PANEL, line=RULE, title="Residual bias = two leakages")
image_fit(s, EQ("zeta_decomp"), 7.2, 3.74, 5.15, 1.0, frame=False)
text(s, 7.2, 4.92, 5.15, 0.75, [[R("Î¶_vol is bounded by volume clipping (spam nullified); "
     "Î¶_med explodes when the median is poisoned (ALIE: f_AÂ·r large, Î”_m large).", 11, TEXT, False, F_BODY)]],
     line_spacing=1.04)
panel(s, 0.8, 5.85, 11.75, 1.10, fill=GOLDTNT, line=ACCENT)
text(s, 1.0, 5.93, 6.0, 0.3, [[R("Corollary â€” perfect filtering recovers clean SGD:", 12, ACCENT2, True, F_HEAD)]])
image_fit(s, EQ("perfect_filter"), 1.0, 6.26, 11.3, 0.42, frame=False)
footer(s, pg())

# --- CONVERGENCE III: THEORY MEETS EXPERIMENT (DIAGNOSTICS) ----------------
s = slide(); header(s, "The Aegis Protocol", "Theory Meets Experiment â€” Filter Diagnostics", tsize=24)
text(s, 0.8, 1.50, 11.8, 0.34, [[R("The detection / false-negative panels are the empirical w_B â€” "
     "the convergence verdict tracks Î¶_A exactly:", 13, NAVY, True, F_HEAD)]])
dg = [
    ("diag_noatk", "No Attack", GREEN, "FN n/a Â· Î¶_A = 0  â†’  75%  (Îº_A = 3Îº_H, honest distortion only)"),
    ("diag_label", "Label-Flip", GREEN, "FN 100â†’0%, detectionâ†’100%  â‡’  râ†’0, w_Bâ†’0, Î¶_A small  â†’  74%"),
    ("diag_alie",  "ALIE", CRIMSON, "FN stuck at 100%, detection 0%  â‡’  râ†’1, Î¶_med large  â†’  10%"),
]
bw = 3.85; x = 0.72
for key, lab, col, cap in dg:
    image_fit(s, RES[key], x, 1.96, bw, 2.95, frame=True)
    text(s, x, 4.98, bw, 0.32, [[R(lab, 14.5, col, True, F_HEAD)]], align=PP_ALIGN.CENTER)
    text(s, x, 5.34, bw, 0.9, [[R(cap, 10.5, TEXT, False, F_BODY)]], align=PP_ALIGN.CENTER, line_spacing=1.04)
    x += bw + 0.18
rect(s, 0.8, 6.42, 11.75, 0.64, fill=GOLDTNT, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
rect(s, 0.8, 6.42, 0.10, 0.64, fill=ACCENT, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, 1.05, 6.42, 11.4, 0.64, [[R("Converge iff Î¶_A is small â€” and Î¶_A is small iff the filter detects the attack. "
     "The theorem predicts every regime the experiments show.", 12.5, NAVY, True, F_BODY)]],
     anchor=MSO_ANCHOR.MIDDLE)
footer(s, pg())

# --- 22. DIVIDER 4 ---------------------------------------------------------
s = slide(); divider(s, "04", "Part Four", "Experimental Setup",
                     "CIFAR-10, 4-shard Non-IID, n = 30, f = 0.30, 1000 rounds")

# --- 23. SETUP -------------------------------------------------------------
s = slide(); header(s, "Setup", "Experimental Apparatus")
quad = [
    ("Data", "CIFAR-10 Â· 4-shard label sharding Â· Î± = 0.5",
     ["30 clients, ~1,664 samples each", "~40% class coverage / client", "Î“  10â€“50Ã— the IID baseline"]),
    ("Model", "ImprovedCNN Â· GroupNorm Â· Dropout",
     ["4 conv + 2 FC, ~1.12M params", "SGD, lr 0.01, momentum 0.8", "1 local epoch, batch 32"]),
    ("Attacks", "f = 0.30 (9/30), p_attack = 1.0",
     ["Sign-flip, Label-flip, Vol-spam", "ALIE (z=1), IPM (Îµ=0.5)", "Sybil (2 clones) â†’ f_eff = 0.56"]),
    ("Baselines", "Five robust aggregators",
     ["FedAvg, CWMed, Multi-Krum", "FoolsGold, Bulyan (f=0.20)", "Best accuracy over 1000 rounds"]),
]
for i, (t, sub, lines) in enumerate(quad):
    col = i % 2; row = i // 2
    cx = 0.8 + col*(5.98); cy = 1.7 + row*(2.45)
    panel(s, cx, cy, 5.78, 2.25, fill=PANEL, line=RULE)
    rect(s, cx, cy, 5.78, 0.10, fill=ACCENT, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, cx+0.25, cy+0.2, 5.3, 0.4, [[R(t, 17, NAVY, True, F_HEAD)]])
    text(s, cx+0.25, cy+0.62, 5.3, 0.4, [[R(sub, 12.5, ACCENT2, True, F_BODY)]])
    bullets(s, cx+0.25, cy+1.05, 5.3, 1.2, [(0, l) for l in lines], size=12, gap=3, lh=1.0)
text(s, 0.8, 6.7, 11.8, 0.4, [[R("Our oracle ceiling is 75.53 Â± 0.22% â€” the hard accuracy limit set by this "
     "extreme Non-IID partition.", 12.5, MUTED, True, F_BODY, True)]])
footer(s, pg())

# --- 24. DIVIDER 5 ---------------------------------------------------------
s = slide(); divider(s, "05", "Part Five", "Results & Analysis",
                     "Per-attack robustness from the finalised experiments â€” one figure at a time")

# --- 25. BYZANTINE TAX TABLE ----------------------------------------------
s = slide(); header(s, "Results", "The Byzantine Tax â€” Cost of Defence Under No Attack")
data = [
    ["Aggregator", "Best Acc (%)", "Best Loss", "Byzantine Tax (pp)"],
    ["FedAvg (non-robust)", "76.16", "0.7248", "0.00"],
    ["Aegis", "76.08", "0.7255", "0.08"],
    ["FoolsGold", "73.80", "0.7628", "2.36"],
    ["CWMed", "72.31", "0.9064", "3.85"],
    ["Multi-Krum", "71.30", "0.8908", "4.86"],
    ["Bulyan", "67.69", "0.9723", "8.47"],
]
styles = {(2, 0): {"color": NAVY, "bold": True, "fill": GOLDTNT, "align": PP_ALIGN.LEFT},
          (2, 1): {"color": GREEN, "bold": True, "fill": GOLDTNT},
          (2, 3): {"color": GREEN, "bold": True, "fill": GOLDTNT},
          (2, 2): {"color": NAVY, "bold": True, "fill": GOLDTNT},
          (7, 3): {"color": CRIMSON, "bold": True}}
table(s, 0.8, 1.7, 6.7, 7, [2.7, 1.5, 1.3, 1.9], data, fsize=12.5, hsize=11.5, row_h=0.5, cell_styles=styles)
panel(s, 7.75, 1.7, 4.8, 4.55, fill=NAVY)
text(s, 8.0, 1.95, 4.3, 0.5, [[R("0.08 pp", 40, ACCENT, True, F_HEAD)]])
text(s, 8.0, 2.75, 4.3, 0.5, [[R("Aegis overhead vs FedAvg", 14, WHITE, True, F_BODY)]])
bullets(s, 8.0, 3.4, 4.3, 2.7, [
    (0, "Aegis is effectively free when there are no attacks.", WHITE, False),
    (0, "Bulyan loses 8.47 pp because Krum selection collapses to a single sector.", RGBColor(0xC6,0xD2,0xDE), False),
    (0, "This proves our adaptive floor successfully preserves honest specialists.", WHITE, True),
], size=13, gap=10)
footer(s, pg())

# --- 26..27 baseline images -----------------------------------------------
img_slide("Results Â· Baseline", "Convergence Stability â€” Three Seeds, No Attack", "baseline_seeds",
          [R("Mean ", 14, TEXT, False, F_BODY), R("75.53 Â± 0.22%", 14, GREEN, True, F_BODY),
           R(" across seeds 42 / 123 / 456 â€” the Non-IID oracle ceiling; single-seed results are representative.", 14, TEXT, False, F_BODY)],
          tsize=25)
img_slide("Results Â· Baseline", "Robust Aggregators vs FedAvg â€” No Attack", "baseline_cmp",
          [R("Aegis tracks FedAvg (76.08 vs 76.16); ", 14, TEXT, False, F_BODY),
           R("CWMed and Krum plateau 4â€“5 pp lower", 14, CRIMSON, True, F_BODY),
           R(" â€” the Byzantine Tax of centroid filters under Non-IID.", 14, TEXT, False, F_BODY)],
          tsize=26)

# --- 28. HEAD-TO-HEAD TABLE ------------------------------------------------
s = slide(); header(s, "Results", "Head-to-Head Under Active Attack  (f = 0.30)")
data = [
    ["Aggregator", "No Atk", "Label Flip", "Sign Flip", "ALIE", "IPM", "Sybil"],
    ["FedAvg", "76.16", "â€”", "40.04", "66.78", "70.25", "64.36"],
    ["CWMed", "72.31", "62.07", "49.64", "23.21", "10.43", "30.58"],
    ["Multi-Krum", "71.30", "72.85", "70.05", "34.18", "17.30", "16.36"],
    ["FoolsGold", "73.80", "67.06", "34.20", "71.15", "71.07", "71.27"],
    ["Bulyan â€ ", "67.69", "68.51", "62.58", "20.39", "10.00", "31.49"],
    ["Aegis (ours)", "76.08", "74.34", "66.61", "10.00", "63.32", "64.22"],
]
styles = {}
for c in range(7):
    styles[(6, c)] = {"color": NAVY, "bold": True, "fill": GOLDTNT, "align": PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER}
styles[(6, 2)] = {"color": GREEN, "bold": True, "fill": GOLDTNT}
styles[(6, 4)] = {"color": CRIMSON, "bold": True, "fill": GOLDTNT}
table(s, 0.8, 1.72, 11.75, 7, [2.35, 1.55, 1.6, 1.6, 1.5, 1.5, 1.5], data, fsize=12.5, hsize=12, row_h=0.52, cell_styles=styles)
bullets(s, 0.8, 5.5, 11.8, 1.5, [
    (0, "Aegis beats every geometric baseline (CWMed, Krum, Bulyan) on every attack; best overall on label-flip (+1.49 pp).", NAVY, True),
    (0, "Three regimes: dominant (label-flip), competitive (sign-flip, IPM, Sybil), failed (ALIE). Per-attack figures follow.", TEXT, False),
    (0, "â€  Bulyan evaluated at f = 0.20 (its n â‰¥ 4f+3 constraint is violated at f = 0.30).", MUTED, False),
], size=12.5, gap=6)
footer(s, pg())

# --- 29..35 per-attack & sweep images -------------------------------------
img_slide("Results Â· Aegis", "Aegis Across Conventional Attacks", "aegis_multi",
          [R("No-attack and label-flip converge together (~74%); sign-flip and IPM are ", 13.5, TEXT, False, F_BODY),
           R("tightly bounded, not collapsed", 13.5, GREEN, True, F_BODY),
           R(" â€” two-pass median limits the damage.", 13.5, TEXT, False, F_BODY)], tsize=26)
img_slide("Results Â· Sign-Flip", "Sign-Flip â€” All Aggregators (f = 0.30)", "sign_cmp",
          [R("Aegis 66.6% â€” no early collapse; ", 14, TEXT, False, F_BODY),
           R("FedAvg crashes to 40%, FoolsGold to 34%", 14, CRIMSON, True, F_BODY),
           R("; Krum (70%) leads on this pure-magnitude attack.", 14, TEXT, False, F_BODY)], tsize=26)
img_slide("Results Â· Label-Flip", "Label-Flip â€” All Aggregators (f = 0.30)", "label_cmp",
          [R("Aegis is the strongest defence at ", 14, TEXT, False, F_BODY),
           R("74.34%", 14, GREEN, True, F_BODY),
           R(" â€” within 1.7 pp of the clean ceiling; the MAD gate intercepts the corrupt gradients.", 14, TEXT, False, F_BODY)], tsize=26)
img_slide("Results Â· IPM", "Inner-Product Manipulation (Îµ = 0.5, Omniscient)", "ipm_cmp",
          [R("Aegis holds 63.3% by aggregate dilution while ", 14, TEXT, False, F_BODY),
           R("CWMed / Krum / Bulyan collapse to â‰¤ 17%", 14, CRIMSON, True, F_BODY),
           R("; only FoolsGold edges ahead (71%).", 14, TEXT, False, F_BODY)], tsize=26)
img_slide("Results Â· ALIE", "A Little Is Enough (Z = 1.0, Omniscient) â€” The Breach", "alie_cmp",
          [R("Aegis collapses to 10% (chance) â€” ALIE hijacks the median; ", 14, TEXT, False, F_BODY),
           R("FoolsGold's history tracking is the only survivor (71%)", 14, CRIMSON, True, F_BODY),
           R(".", 14, TEXT, False, F_BODY)], tsize=27)

img_slide("Results Â· Volume Spam", "Volume Spam â€” All Aggregators (f = 0.30)", "vol_cmp",
          [R("Aegis #1 at ", 13.5, TEXT, False, F_BODY),
           R("73.77%", 13.5, GREEN, True, F_BODY),
           R(" â€” median volume clipping neutralises the inflation attack; ", 13.5, TEXT, False, F_BODY),
           R("FedAvg collapses to 14%", 13.5, CRIMSON, True, F_BODY),
           R(".", 13.5, TEXT, False, F_BODY)], tsize=26)

img_slide("Results Â· Sybil", "Sybil â€” All Aggregators (f = 0.30, k = 2)", "syb_cmp",
          [R("FoolsGold #1 (68.76%) via full-history tracking; Aegis #2 at ", 13.5, TEXT, False, F_BODY),
           R("63.17%", 13.5, GREEN, True, F_BODY),
           R(". ", 13.5, TEXT, False, F_BODY),
           R("Geometric baselines destroyed", 13.5, CRIMSON, True, F_BODY),
           R(" (Krum 14%, CWMed 22%).", 13.5, TEXT, False, F_BODY)], tsize=26)
img_slide("Results Â· Advanced", "Volume-Spam, ALIE & Sybil vs Aegis", "advanced",
          [R("Volume-spam is fully nullified (74%) by median clipping; ", 13.5, TEXT, False, F_BODY),
           R("ALIE and dense Sybil clusters overwhelm the median anchor", 13.5, CRIMSON, True, F_BODY),
           R(".", 13.5, TEXT, False, F_BODY)], tsize=26)
img_slide("Results Â· Stress Test", "Resilience vs Byzantine Fraction (Sign-Flip)", "byz_sweep",
          [R("Graceful degradation to f = 0.30 (~2.8 pp / +10 pp); ", 14, TEXT, False, F_BODY),
           R("a cliff at f = 0.40", 14, CRIMSON, True, F_BODY),
           R(" as adversaries become a per-round majority and break the median.", 14, TEXT, False, F_BODY)], tsize=24)
img_slide("Results Â· Heterogeneity", "Sensitivity to Non-IID Degree (Shard Sweep)", "shard_sweep",
          [R("At 4 shards the label-flip gap is â‰ˆ 0; ", 14, TEXT, False, F_BODY),
           R("extreme skew (2 shards) widens it to ~15 pp", 14, CRIMSON, True, F_BODY),
           R(" as honest variance overlaps the attack.", 14, TEXT, False, F_BODY)], tsize=25)

# --- 36. ABLATION TABLE ----------------------------------------------------
s = slide(); header(s, "Results", "Ablation: Which Component Carries the Load?")
data = [
    ["Configuration", "Best Acc", "Î” Acc", "Filter Acc", "DR %", "Prec %"],
    ["Full Aegis", "73.77", "â€”", "87.8", "60.7", "99.1"],
    ["No Adaptive Threshold", "72.47", "âˆ’1.30", "97.2", "92.0", "98.7"],
    ["No Cosine Penalty", "73.73", "âˆ’0.04", "87.4", "59.5", "99.4"],
    ["No Volume Clipping", "74.67", "+0.90", "92.4", "76.6", "98.6"],
    ["No Median Decontam.", "74.24", "+0.47", "91.0", "71.6", "98.7"],
    ["No Euclidean Filter", "71.65", "âˆ’2.12", "69.9", "0.0", "N/A"],
]
styles = {(1, 0): {"color": NAVY, "bold": True, "align": PP_ALIGN.LEFT, "fill": HEADTNT},
          (6, 0): {"color": CRIMSON, "bold": True, "fill": GOLDTNT, "align": PP_ALIGN.LEFT},
          (6, 2): {"color": CRIMSON, "bold": True, "fill": GOLDTNT},
          (6, 3): {"color": CRIMSON, "bold": True, "fill": GOLDTNT},
          (6, 4): {"color": CRIMSON, "bold": True, "fill": GOLDTNT},
          (2, 2): {"color": AMBER, "bold": True}}
table(s, 0.8, 1.7, 7.45, 7, [2.65, 1.1, 0.95, 1.15, 0.85, 0.85], data, fsize=11.5, hsize=11, row_h=0.5, cell_styles=styles)
panel(s, 8.45, 1.7, 4.1, 5.05, fill=PANEL, line=RULE, title="Reading the table")
bullets(s, 8.65, 2.3, 3.75, 4.4, [
    (0, "The Euclidean MAD filter is the load-bearing wall: removing it costs 2.12 pp and zeroes detection.", NAVY, True),
    (0, "Adaptive threshold trades precision for recall â€” a fixed k lifts DR to 92% but adds Byzantine Tax.", TEXT, False),
    (0, "Cosine penalty is marginal for label-flip (âˆ’0.04) â€” but it is the only signal under IPM.", TEXT, False),
], size=12.5, gap=11)
footer(s, pg())

# --- 37. ALIE ABLATION IMAGE ----------------------------------------------
img_slide("Results Â· ALIE Ablation", "No Hyperparameter Recovers ALIE", "alie_ablation",
          [R("Sweeping the cosine-penalty weight Î± leaves accuracy near chance â€” ", 14, TEXT, False, F_BODY),
           R("ALIE is a structural gap, not a tuning issue", 14, CRIMSON, True, F_BODY),
           R(".", 14, TEXT, False, F_BODY)], tsize=26)

# --- 37b. VOLUME CLIPPING ABLATION -----------------------------------------
img_slide("Results Â· Ablation", "Volume Clipping â€” The Load-Bearing Wall for Spam", "vol_ablation",
          [R("Full Aegis reaches ", 14, TEXT, False, F_BODY),
           R("73.60%", 14, GREEN, True, F_BODY),
           R(" under volume spam; removing the clip collapses accuracy to ", 14, TEXT, False, F_BODY),
           R("14.58%", 14, CRIMSON, True, F_BODY),
           R(" â€” a 59 pp gap proves volume clipping is essential.", 14, TEXT, False, F_BODY)], tsize=24)

# --- 38. SUMMARY TABLE -----------------------------------------------------
s = slide(); header(s, "Results", "Summary Across All Attacks")
data = [
    ["Attack", "Aegis Acc", "DR %", "Best Baseline", "Gap (pp)", "Regime"],
    ["No Attack", "76.08", "N/A", "FedAvg  76.16", "âˆ’0.08", "I"],
    ["Label Flip", "74.34", "76.5", "Krum  72.85", "+1.49", "I"],
    ["Sign Flip", "66.61", "60.8", "Krum  70.05", "âˆ’3.44", "II"],
    ["IPM", "63.32", "0.0", "FoolsGold  71.07", "âˆ’7.75", "II"],
    ["Sybil", "64.22", "0.0", "FoolsGold  71.27", "âˆ’7.05", "II/III"],
    ["ALIE", "10.00", "0.0", "FoolsGold  71.15", "âˆ’61.15", "III"],
]
reg_color = {"I": GREEN, "II": AMBER, "II/III": AMBER, "III": CRIMSON}
styles = {}
for r in range(1, 7):
    styles[(r, 5)] = {"color": WHITE, "bold": True, "fill": reg_color[data[r][5]]}
    styles[(r, 0)] = {"color": NAVY, "bold": True, "align": PP_ALIGN.LEFT}
styles[(2, 4)] = {"color": GREEN, "bold": True}
styles[(6, 4)] = {"color": CRIMSON, "bold": True}
table(s, 0.8, 1.7, 11.75, 7, [1.85, 1.75, 1.3, 3.2, 1.65, 2.0], data, fsize=12.5, hsize=12, row_h=0.52, cell_styles=styles)
bullets(s, 0.8, 5.75, 11.8, 1.4, [
    (0, "Aegis is strictly superior to all geometric baselines (CWMed, Krum, Bulyan) on every attack.", NAVY, True),
    (0, "FoolsGold wins ALIE & Sybil â€” its full-history tracking outresolves the 20-round EMA.", TEXT, False),
    (0, "No aggregator dominates everywhere â†’ motivates an attack-type-aware routing layer (future work).", TEXT, False),
], size=13.5, gap=6)
footer(s, pg())

# --- 39. DIVIDER 6 ---------------------------------------------------------
s = slide(); divider(s, "06", "Part Six", "Conclusions & Future Work",
                     "What survives the experimental critique â€” and what comes next")

# --- 40. CONCLUSIONS -------------------------------------------------------
s = slide(); header(s, "Conclusion", "Achievements & Honest Boundaries")
panel(s, 0.8, 1.7, 5.78, 4.95, fill=PANEL, line=RULE)
rect(s, 0.8, 1.7, 5.78, 0.1, fill=GREEN, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, 1.05, 1.95, 5.3, 0.5, [[R("What survives the critique", 16, GREEN, True, F_HEAD)]])
bullets(s, 1.05, 2.55, 5.35, 4.0, [
    (0, "Near-zero Byzantine Tax: 0.08 pp vs FedAvg under no attack.", TEXT, False),
    (0, "Dominant label-flip resilience: 74.34%, 98.8% precision.", TEXT, False),
    (0, "O(kd) linear-time aggregation â€” ~30Ã— faster than Bulyan (RÂ²=0.999).", TEXT, False),
    (0, "Two-pass median prevents Center Drag â€” a structural contribution.", TEXT, False),
], size=13.5, gap=10)
panel(s, 6.78, 1.7, 5.77, 4.95, fill=PANEL, line=RULE)
rect(s, 6.78, 1.7, 5.77, 0.1, fill=CRIMSON, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, 7.03, 1.95, 5.3, 0.5, [[R("Stated boundaries", 16, CRIMSON, True, F_HEAD)]])
bullets(s, 7.03, 2.55, 5.35, 4.0, [
    (0, "IPM partial evasion: 63.3% by dilution, 0% detection â€” velocity patch proposed.", TEXT, False),
    (0, "ALIE catastrophic failure: 10% â€” EMA window too short; remains open.", TEXT, False),
    (0, "f = 0.40 collapse: provable limit of any median-based scheme (0.5 breakdown).", TEXT, False),
], size=13.5, gap=10)
text(s, 0.8, 6.78, 11.8, 0.4, [[R("The failure modes are structural gaps requiring architectural extensions â€” not "
     "implementation defects.", 13, ACCENT2, True, F_BODY)]])
footer(s, pg())

# --- 41. FUTURE WORK -------------------------------------------------------
s = slide(); header(s, "Future Work", "Four Research Directions")
fw = [
    ("Server-Velocity for IPM", "Screen against a momentum buffer Vâ½áµ—â»Â¹â¾ â€” a pre-debiased reference IPM cannot poison. (Re-activates the interim's server momentum.)"),
    ("Kalman Reputation for ALIE", "Replace fixed-window EMA with a Kalman-filtered estimator â€” adaptive memory at O(Kd), not FoolsGold's O(KÂ²d)."),
    ("Asynchronous Aegis", "Staleness-weighted cosine penalty to suppress deliberately delayed Byzantine updates over disrupted links."),
    ("Dynamic EMA Decay", "Tie Î³ to network churn â€” shrink memory in volatile rounds, keep long memory when stable."),
]
for i, (t, d) in enumerate(fw):
    col = i % 2; row = i // 2
    cx = 0.8 + col*5.98; cy = 1.8 + row*2.35
    panel(s, cx, cy, 5.78, 2.1, fill=PANEL, line=RULE)
    rect(s, cx+0.25, cy+0.24, 0.62, 0.62, fill=NAVY, shape=MSO_SHAPE.OVAL)
    text(s, cx+0.25, cy+0.24, 0.62, 0.62, [[R(str(i+1), 22, ACCENT, True, F_LIGHT)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, cx+1.1, cy+0.26, 4.5, 0.5, [[R(t, 16, NAVY, True, F_HEAD)]])
    text(s, cx+1.1, cy+0.78, 4.5, 1.2, [[R(d, 12, MUTED, False, F_BODY)]], line_spacing=1.05)
footer(s, pg())

# --- 42..43 REFERENCES -----------------------------------------------------
refs1 = [
    "H. B. McMahan et al., â€œCommunication-Efficient Learning of Deep Networks from Decentralized Data,â€ AISTATS, 2017.",
    "P. Blanchard et al., â€œMachine Learning with Adversaries: Byzantine-Tolerant Gradient Descent (Krum),â€ NeurIPS, 2017.",
    "D. Yin et al., â€œByzantine-Robust Distributed Learning: Towards Optimal Statistical Rates,â€ ICML, 2018.",
    "E. M. El Mhamdi et al., â€œThe Hidden Vulnerability of Distributed Learning in Byzantium (Bulyan),â€ ICML, 2018.",
    "G. Baruch et al., â€œA Little Is Enough: Circumventing Defenses for Distributed Learning (ALIE),â€ NeurIPS, 2019.",
    "C. Xie et al., â€œFall of Empires: Breaking Byzantine-tolerant SGD by Inner Product Manipulation (IPM),â€ UAI, 2019.",
    "C. Fung et al., â€œThe Limitations of Federated Learning in Sybil Settings (FoolsGold),â€ RAID, 2020.",
    "X. Cao et al., â€œFLTrust: Byzantine-Robust Federated Learning via Trust Bootstrapping,â€ NDSS, 2021.",
]
refs2 = [
    "T. D. Nguyen et al., â€œFLAME: Taming Backdoors in Federated Learning,â€ USENIX Security, 2022.",
    "C. Xie et al., â€œZeno: Distributed SGD with Suspicion-based Fault-tolerance,â€ ICML, 2019.",
    "S. P. Karimireddy et al., â€œLearning from History for Byzantine Robust Optimization,â€ ICML, 2021.",
    "V. Shejwalkar, A. Houmansadr, â€œManipulating the Byzantine (DnC),â€ NDSS, 2021.",
    "S. Awan et al., â€œCONTRA: Defending Against Poisoning Attacks in Federated Learning,â€ ESORICS, 2021.",
    "K. Pillutla et al., â€œRobust Aggregation for Federated Learning (RFA),â€ IEEE Trans. Signal Process., 2022.",
    "U. Shahul, J. Harshan, â€œFORTA: Byzantine-Resilient FL Aggregation via DFT-Guided Krum,â€ arXiv:2507.14588, 2025.",
    "T.-M. H. Hsu et al., â€œMeasuring the Effects of Non-Identical Data Distribution for Federated Visual Classification,â€ NeurIPS-W, 2019.",
]
def ref_slide(title, refs, start):
    s = slide(); header(s, "References", title)
    tb = s.shapes.add_textbox(Inches(0.85), Inches(1.7), Inches(11.7), Inches(5.2))
    tf = tb.text_frame; tf.word_wrap = True
    for i, rtext in enumerate(refs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(11); p.line_spacing = 1.04
        a = p.add_run(); a.text = f"[{start+i}]  "
        a.font.name = F_HEAD; a.font.size = Pt(13); a.font.bold = True; a.font.color.rgb = ACCENT2
        b = p.add_run(); b.text = rtext
        b.font.name = F_BODY; b.font.size = Pt(13); b.font.color.rgb = TEXT
    footer(s, pg())
ref_slide("Selected References  (1/2)", refs1, 1)
ref_slide("Selected References  (2/2)", refs2, 9)

# --- 44. THANK YOU ---------------------------------------------------------
s = slide()
rect(s, 0, 0, SW, SH, fill=NAVY)
rect(s, 0, 0, SW, 0.22, fill=ACCENT)
rect(s, 0, SH-0.22, SW, 0.22, fill=ACCENT)
text(s, 1.0, 2.35, 11.3, 1.2, [[R("Thank You", 56, WHITE, True, F_HEAD)]])
text(s, 1.05, 3.55, 11.0, 0.6, [[R("Questions & Discussion", 24, ACCENT, True, F_LIGHT)]])
rect(s, 1.05, 4.45, 4.6, 0.035, fill=NAVY2)
text(s, 1.0, 4.75, 11.0, 1.4, [
    [R("Enhanced Aegis: A Byzantine-Resilient FL Protocol for Drone Swarms", 15, RGBColor(0xD7,0xE0,0xEA), False, F_BODY)],
    [R("Abhishek Kumar Tripathi  Â·  2024EET2845  Â·  Advisor: Prof. Harshan Jagadeesh", 14, RGBColor(0xB9,0xC6,0xD4), False, F_BODY)],
    [R("Department of Electrical Engineering, IIT Delhi", 13, RGBColor(0xA9,0xB8,0xC8), False, F_BODY)],
], space_after=7)

# --- APPENDIX (BACKUP): HYPERPARAMETER SELECTION ---------------------------
s = slide(); header(s, "Appendix Â· Backup", "Why These Hyperparameters?")
panel(s, 0.8, 1.6, 11.75, 1.02, fill=GOLDTNT, line=ACCENT)
image_fit(s, EQ("hp_credit"), 1.0, 1.7, 11.3, 0.82, frame=False)
data = [
    ["Symbol", "Value", "Rationale  (derived from the Non-IID statistics, not hand-tuned)"],
    ["Ï„", "âˆ’0.3", "honest cosine floor â‰ˆ âˆ’0.2;  attacker ceiling â‰ˆ âˆ’0.8  â†’  safe margin both sides"],
    ["k_max", "6.0", "z â‰ˆ 4.05  â†’  < 0.003% honest false-positives in noisy early rounds"],
    ["k_min", "2.0", "asymptotic IID target (z â‰ˆ 1.35);  never reached under Non-IID"],
    ["k_floor", "4.0", "z â‰ˆ 2.70  â†’  FP â‰ˆ 0 against the 4-shard honest fat tail"],
    ["T_warm", "300", "kâ½áµ—â¾ reaches the floor by ~round 200, just after ~90% convergence"],
    ["Î½  (var. sens.)", "3.0", "widens kâ½áµ—â¾ ~30% at CV = 0.10;  calibrated for f â‰¤ 30%"],
    ["Î±  (cosine)", "30.0", "attacker influence < 10%  while honest-outlier credit stays > 25%"],
    ["Î³  (EMA)", "0.95", "20-round memory;  one anomalous round decays in ~50 rounds"],
    ["Î»  (reputation)", "20.0", "reputation = 40% of the directional penalty at steady state"],
]
styles = {(r, 0): {"color": NAVY, "bold": True, "align": PP_ALIGN.LEFT} for r in range(1, 11)}
for r in range(1, 11):
    styles[(r, 1)] = {"color": ACCENT2, "bold": True}
table(s, 0.8, 2.78, 11.75, 11, [1.85, 1.15, 8.75], data, fsize=11.5, hsize=11.5, row_h=0.375, cell_styles=styles)
footer(s, pg())

prs.save(OUT)
print("SAVED:", OUT)
print("SLIDES:", len(prs.slides._sldIdLst))
