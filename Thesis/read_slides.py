import sys
from pptx import Presentation

prs = Presentation(r'd:\IITD\MTP 2\FL_Project\Thesis\Aegis_Defense.pptx')
with open('slides.txt', 'w', encoding='utf-8') as f:
    for i, slide in enumerate(prs.slides):
        f.write(f'--- Slide {i+1} ---\n')
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                f.write(shape.text + '\n')
