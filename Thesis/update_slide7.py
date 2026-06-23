import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR

NAVY    = RGBColor(0x0E, 0x2A, 0x47)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT   = RGBColor(0xF5, 0xF7, 0xFA)
TEXT    = RGBColor(0x1F, 0x29, 0x33)
CRIMSON = RGBColor(0xC0, 0x39, 0x2B)

F_HEAD = "Segoe UI Semibold"
F_BODY = "Segoe UI"

def set_cell(cell, txt, size, color, bold=False, align=PP_ALIGN.CENTER, fill=WHITE, font=F_BODY, italic=False):
    cell.fill.solid(); cell.fill.fore_color.rgb = fill
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = Inches(0.06); cell.margin_right = Inches(0.06)
    cell.margin_top = Inches(0.02); cell.margin_bottom = Inches(0.02)
    tf = cell.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = txt
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color; r.font.name = font

def table(s, x, y, w, rows, col_w, data, header_fill=NAVY, header_color=WHITE,
          fsize=10.5, hsize=10.5, row_h=0.4, cell_styles=None, zebra=True):
    nr = len(data); nc = len(data[0])
    gt = s.shapes.add_table(nr, nc, Inches(x), Inches(y), Inches(w), Inches(row_h*nr)).table
    gt.first_row = False; gt.horz_banding = False
    for ci, cw in enumerate(col_w):
        gt.columns[ci].width = Inches(cw)
    for ri in range(nr):
        gt.rows[ri].height = Inches(row_h)
        for ci in range(nc):
            cell = gt.cell(ri, ci)
            txt = str(data[ri][ci])
            if ri == 0:
                set_cell(cell, txt, hsize, header_color, True,
                         PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER, header_fill, F_HEAD)
            else:
                base = LIGHT if (zebra and ri % 2 == 0) else WHITE
                col = TEXT; bold = False; al = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
                st = (cell_styles or {}).get((ri, ci))
                if st:
                    col = st.get("color", col); bold = st.get("bold", bold); base = st.get("fill", base)
                    al = st.get("align", al)
                set_cell(cell, txt, fsize, col, bold, al, base,
                         F_HEAD if (ci == 0 and bold) else F_BODY,
                         italic=(st or {}).get("italic", False) if st else False)
    return gt

prs = Presentation(r'd:\IITD\MTP 2\FL_Project\Thesis\Aegis_Defense.pptx')
slide = prs.slides[8]

# Delete old table
for shape in list(slide.shapes):
    if shape.has_table:
        sp = shape._element
        sp.getparent().remove(sp)

data = [
    ["Attack", "Mechanism", "Norm\nevasion", "Dir.\nevasion", "Median\npoison", "Identity\nspoof"],
    ["Label-Flip", "Data corruption", "✓", "✓", "✕", "✕"],
    ["Sign-Flip", "Direction inversion", "✕", "✕", "✕", "✕"],
    ["IPM", "Stealth anti-gradient", "✓", "✕", "✓", "✕"],
    ["ALIE", "In-variance poison", "✓", "✓", "✓", "✕"],
    ["Sybil", "Density flooding", "✓", "✓", "✓", "✓"],
    ["Volume Spam", "Dataset inflation", "✕", "✕", "✕", "✕"],
]
styles = {}
for r in range(1, 7):
    for c in range(2, 6):
        v = data[r][c]
        styles[(r, c)] = {"color": CRIMSON if v == "✓" else RGBColor(0x9A,0xA6,0xB2), "bold": v == "✓"}

table(slide, 0.8, 3.55, 11.75, 7, [1.85, 3.5, 1.6, 1.6, 1.6, 1.6], data,
      fsize=12.5, hsize=11, row_h=0.45, cell_styles=styles)

prs.save(r'd:\IITD\MTP 2\FL_Project\Thesis\Aegis_Defense.pptx')
