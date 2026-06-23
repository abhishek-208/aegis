import os
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

def image_fit(s, path, bx, by, bw, bh):
    iw, ih = Image.open(path).size
    ar = iw/ih; bar = bw/bh
    if ar > bar:
        w = bw; h = bw/ar
    else:
        h = bh; w = bh*ar
    x = bx + (bw-w)/2; y = by + (bh-h)/2
    s.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))

prs = Presentation(r'd:\IITD\MTP 2\FL_Project\Thesis\Aegis_Defense.pptx')

removed_footers = 0
for slide in prs.slides:
    for shape in list(slide.shapes):
        if shape.has_text_frame:
            text = shape.text.replace('\n', '').strip()
            if "Enhanced Aegis" in text and "Thesis Defense" in text:
                sp = shape._element
                sp.getparent().remove(sp)
                removed_footers += 1

print(f"Removed {removed_footers} footer text boxes.")

# Update slide 21 (index 20) with the new complexity image
slide21 = prs.slides[20]
for shape in list(slide21.shapes):
    if getattr(shape, "shape_type", None) == 13: # PICTURE
        # The equation image is the top one (y < 2.0 inches)
        if shape.top < Inches(2.0):
            sp = shape._element
            sp.getparent().remove(sp)
            print("Removed old complexity equation image.")

# Add new complexity image
img_path = r'd:\IITD\MTP 2\FL_Project\Thesis\_math\complexity.png'
image_fit(slide21, img_path, 1.0, 1.74, 11.3, 0.78)
print("Inserted new complexity equation image.")

try:
    prs.save(r'd:\IITD\MTP 2\FL_Project\Thesis\Aegis_Defense.pptx')
    print("Saved successfully.")
except PermissionError:
    print("PermissionError: PowerPoint file is locked.")
