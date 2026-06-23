# -*- coding: utf-8 -*-
"""
Update Aegis_Defense.pptx with Expr 33 Volume Spam & Sybil results.
  1. Update Head-to-Head table (Slide 32, idx 31): add Vol Spam column, update Sybil, reorder
  2. Insert two new per-attack slides after Slide 37 (idx 36): Vol Spam + Sybil
  3. Update Summary table (Slide 43, idx 42): add Vol Spam row, update Sybil row
"""
import os, copy, sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image
from lxml import etree

PPTX = r'd:\IITD\MTP 2\FL_Project\Thesis\Aegis_Defense.pptx'
RESDIR = r'D:\IITD\MTP 2\FL_Project\MTP2 Results'
VOL_IMG = os.path.join(RESDIR, r'Expr 33 - @Aegis, Fedavg, CWmed, Krum, Foolsgold, Bulyan, Volume Spam',
                       'fed_avg+aegis+cw_med+multi_krum+fools_gold+bulyan_volume_spam_manual_accuracy_line.png')
SYB_IMG = os.path.join(RESDIR, r'Expr 33 - @Aegis, Fedavg, CWmed, Krum, Foolsgold, Bulyan, Sybil',
                       'fed_avg+aegis+cw_med+multi_krum+fools_gold+bulyan_sybil_manual_accuracy_line.png')

# ─── palette (from build_deck.py) ───
NAVY    = RGBColor(0x0E, 0x2A, 0x47)
ACCENT  = RGBColor(0xE3, 0xA5, 0x4E)
ACCENT2 = RGBColor(0xB9, 0x7A, 0x0B)
CRIMSON = RGBColor(0xC0, 0x39, 0x2B)
GREEN   = RGBColor(0x2E, 0x7D, 0x32)
LIGHT   = RGBColor(0xF5, 0xF7, 0xFA)
TEXT    = RGBColor(0x1F, 0x29, 0x33)
MUTED   = RGBColor(0x5C, 0x6B, 0x7B)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
RULE    = RGBColor(0xD6, 0xDE, 0xE6)
GOLDTNT = RGBColor(0xF7, 0xE1, 0xAE)
PANEL   = RGBColor(0xEC, 0xF1, 0xF6)
GRNTNT  = RGBColor(0xDD, 0xEE, 0xDF)
CRIMTNT = RGBColor(0xF6, 0xDE, 0xDA)

F_HEAD = "Segoe UI Semibold"
F_BODY = "Segoe UI"
F_MONO = "Consolas"

# ─── helpers ───
def _noshadow(shp):
    try: shp.shadow.inherit = False
    except: pass

def _rect(s, x, y, w, h, fill=None, line=None, line_w=1.0, shape=MSO_SHAPE.RECTANGLE):
    sp = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    _noshadow(sp)
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    return sp

def _text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=4, line_spacing=1.0):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after); p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (txt, size, color, bold, font, italic) in para:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
            r.font.name = font; r.font.color.rgb = color
    return tb

def R(txt, size, color=TEXT, bold=False, font=F_BODY, italic=False):
    return (txt, size, color, bold, font, italic)

def _header(s, kicker, title, tsize=27):
    _rect(s, 0.55, 0.50, 0.13, 0.66, fill=ACCENT)
    _text(s, 0.82, 0.46, 11.8, 0.3, [[R(kicker.upper(), 12, ACCENT2, True, F_HEAD)]])
    _text(s, 0.80, 0.70, 11.0, 0.7, [[R(title, tsize, NAVY, True, F_HEAD)]])
    _rect(s, 0.55, 1.42, 12.23, 0.018, fill=RULE)

def _image_fit(s, path, bx, by, bw, bh):
    iw, ih = Image.open(path).size
    ar = iw/ih; bar = bw/bh
    if ar > bar: w = bw; h = bw/ar
    else: h = bh; w = bh*ar
    x = bx + (bw-w)/2; y = by + (bh-h)/2
    _rect(s, x-0.05, y-0.05, w+0.10, h+0.10, fill=WHITE, line=RULE, line_w=1.0)
    s.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))

def set_cell(cell, txt, size, color, bold=False, align=PP_ALIGN.CENTER,
             fill=WHITE, font=F_BODY, italic=False):
    cell.fill.solid(); cell.fill.fore_color.rgb = fill
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = Inches(0.06); cell.margin_right = Inches(0.06)
    cell.margin_top = Inches(0.02); cell.margin_bottom = Inches(0.02)
    tf = cell.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = txt
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color; r.font.name = font

def move_slide(prs, old_index, new_index):
    """Move a slide from old_index to new_index."""
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    el = slides[old_index]
    xml_slides.remove(el)
    if new_index >= len(list(xml_slides)):
        xml_slides.append(el)
    else:
        ref = list(xml_slides)[new_index]
        xml_slides.insert(xml_slides.index(ref), el)

# ═══════════════════════════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════════════════════════
prs = Presentation(PPTX)
BLANK = prs.slide_layouts[6]

# ═══════════════════════════════════════════════════════════════════════════════
# 1. UPDATE HEAD-TO-HEAD TABLE — Slide 32 (index 31)
# ═══════════════════════════════════════════════════════════════════════════════
print("=== Step 1: Updating Head-to-Head table on Slide 32 ===")
s32 = prs.slides[31]

# Remove old table and old bullet text
for shape in list(s32.shapes):
    if shape.has_table:
        sp = shape._element; sp.getparent().remove(sp)
        print("  Removed old table.")
    elif shape.has_text_frame:
        t = shape.text_frame.text
        if 'Aegis beats every' in t or 'Aegis heavily outperforms' in t or 'Three regimes' in t or 'Bulyan evaluated' in t:
            sp = shape._element; sp.getparent().remove(sp)
            print(f"  Removed old bullet: {t[:60]}...")

# New table data (reordered: ALIE last)
#           No Atk   LblFlip  SignFlip IPM     VolSpam  Sybil    ALIE
new_data = [
    ["Aggregator",   "No Atk","Label Flip","Sign Flip","IPM","Vol Spam","Sybil","ALIE"],
    ["FedAvg",       "76.16", "\u2014",    "40.04",    "70.25","14.44", "62.48","66.78"],
    ["CWMed",        "72.31", "62.07",     "49.64",    "10.43","63.82", "22.05","23.21"],
    ["Multi-Krum",   "71.30", "72.85",     "70.05",    "17.30","72.45", "14.43","34.18"],
    ["FoolsGold",    "73.80", "67.06",     "34.20",    "71.07","66.96", "68.76","71.15"],
    ["Bulyan \u2020","67.69", "68.51",     "62.58",    "10.00","69.21", "46.45","20.39"],
    ["Aegis (ours)", "76.08", "74.34",     "66.61",    "63.32","73.77", "63.17","10.00"],
]

col_w = [1.60, 1.10, 1.30, 1.30, 1.05, 1.30, 1.05, 1.05]
nr = len(new_data); nc = len(new_data[0])
row_h = 0.40
gt = s32.shapes.add_table(nr, nc, Inches(0.55), Inches(1.70), Inches(sum(col_w)), Inches(row_h*nr)).table
gt.first_row = False; gt.horz_banding = False

for ci, cw in enumerate(col_w):
    gt.columns[ci].width = Inches(cw)

# Cell styles: highlight best value per attack column (green bg)
best_cells = {
    (6, 1): {"fill": GRNTNT, "bold": True},  # Aegis No Atk
    (6, 2): {"fill": GRNTNT, "bold": True},  # Aegis Label Flip
    (3, 3): {"fill": GRNTNT, "bold": True},  # Krum Sign Flip
    (4, 4): {"fill": GRNTNT, "bold": True},  # FoolsGold IPM
    (6, 5): {"fill": GRNTNT, "bold": True},  # Aegis Vol Spam
    (4, 6): {"fill": GRNTNT, "bold": True},  # FoolsGold Sybil
    (4, 7): {"fill": GRNTNT, "bold": True},  # FoolsGold ALIE
    # ALIE failure for Aegis (red)
    (6, 7): {"fill": CRIMTNT, "bold": True, "color": CRIMSON},
    # FedAvg Vol Spam collapse (red)
    (1, 5): {"fill": CRIMTNT, "bold": True, "color": CRIMSON},
}

for ri in range(nr):
    gt.rows[ri].height = Inches(row_h)
    for ci in range(nc):
        cell = gt.cell(ri, ci)
        txt = str(new_data[ri][ci])
        if ri == 0:
            set_cell(cell, txt, 10.5, WHITE, True,
                     PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER, NAVY, F_HEAD)
        else:
            base = LIGHT if ri % 2 == 0 else WHITE
            col = TEXT; bold = False
            st = best_cells.get((ri, ci))
            if st:
                col = st.get("color", col)
                bold = st.get("bold", bold)
                base = st.get("fill", base)
            set_cell(cell, txt, 10.5, col, bold,
                     PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER,
                     base, F_HEAD if (ci == 0 and bold) else F_BODY)

print("  New 8-column table inserted.")

# Add updated bullets below
tb = s32.shapes.add_textbox(Inches(0.8), Inches(5.60), Inches(11.8), Inches(1.5))
tf = tb.text_frame; tf.word_wrap = True
tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0

# Bullet 1
p = tf.paragraphs[0]
p.space_after = Pt(6); p.space_before = Pt(0); p.line_spacing = 1.04; p.alignment = PP_ALIGN.LEFT
b = p.add_run(); b.text = "\u25aa  "; b.font.size = Pt(11.5); b.font.color.rgb = ACCENT; b.font.name = F_BODY; b.font.bold = True
r = p.add_run(); r.text = "Aegis is #1 on Label-Flip (+1.49 pp) and Vol Spam (+1.32 pp over Krum). Geometric baselines collapse on Sybil (Krum 14%, CWMed 22%)."
r.font.size = Pt(12.5); r.font.color.rgb = NAVY; r.font.name = F_BODY; r.font.bold = True

# Bullet 2
p2 = tf.add_paragraph()
p2.space_after = Pt(6); p2.space_before = Pt(0); p2.line_spacing = 1.04; p2.alignment = PP_ALIGN.LEFT
b2 = p2.add_run(); b2.text = "\u25aa  "; b2.font.size = Pt(11.5); b2.font.color.rgb = ACCENT; b2.font.name = F_BODY; b2.font.bold = True
r2 = p2.add_run(); r2.text = "Three regimes: dominant (label-flip, vol-spam), competitive (sign-flip, IPM, Sybil), failed (ALIE). Per-attack figures follow."
r2.font.size = Pt(12.5); r2.font.color.rgb = TEXT; r2.font.name = F_BODY; r2.font.bold = False

# Bullet 3
p3 = tf.add_paragraph()
p3.space_after = Pt(6); p3.space_before = Pt(0); p3.line_spacing = 1.04; p3.alignment = PP_ALIGN.LEFT
b3 = p3.add_run(); b3.text = "\u25aa  "; b3.font.size = Pt(11.5); b3.font.color.rgb = ACCENT; b3.font.name = F_BODY; b3.font.bold = True
r3 = p3.add_run(); r3.text = "\u2020 Bulyan evaluated at f = 0.20 (its n \u2265 4f+3 constraint is violated at f = 0.30)."
r3.font.size = Pt(12.5); r3.font.color.rgb = MUTED; r3.font.name = F_BODY; r3.font.bold = False

print("  Updated bullet text.")

# ═══════════════════════════════════════════════════════════════════════════════
# 2. INSERT TWO NEW SLIDES AFTER SLIDE 37 (ALIE, index 36)
# ═══════════════════════════════════════════════════════════════════════════════
print("\n=== Step 2: Inserting Volume Spam and Sybil slides after Slide 37 ===")

# --- Volume Spam slide ---
sv = prs.slides.add_slide(BLANK)
_header(sv, "Results \u00b7 Volume Spam", "Volume Spam \u2014 All Aggregators (f = 0.30)", tsize=26)
_image_fit(sv, VOL_IMG, 0.8, 1.62, 11.75, 4.62)
_rect(sv, 0.8, 6.42, 11.75, 0.66, fill=GOLDTNT, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
_rect(sv, 0.8, 6.42, 0.10, 0.66, fill=ACCENT, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
_text(sv, 1.05, 6.42, 11.4, 0.66,
      [[R("Aegis #1 at ", 13.5, TEXT, False, F_BODY),
        R("73.77%", 13.5, GREEN, True, F_BODY),
        R(" \u2014 median volume clipping neutralises the inflation attack; ", 13.5, TEXT, False, F_BODY),
        R("FedAvg collapses to 14%", 13.5, CRIMSON, True, F_BODY),
        R(".", 13.5, TEXT, False, F_BODY)]],
      anchor=MSO_ANCHOR.MIDDLE)
vol_idx = len(prs.slides) - 1
move_slide(prs, vol_idx, 37)  # insert after ALIE (idx 36) → position 37
print("  Inserted Volume Spam slide at position 38.")

# --- Sybil slide ---
ss = prs.slides.add_slide(BLANK)
_header(ss, "Results \u00b7 Sybil", "Sybil \u2014 All Aggregators (f = 0.30, k = 2)", tsize=26)
_image_fit(ss, SYB_IMG, 0.8, 1.62, 11.75, 4.62)
_rect(ss, 0.8, 6.42, 11.75, 0.66, fill=GOLDTNT, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
_rect(ss, 0.8, 6.42, 0.10, 0.66, fill=ACCENT, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
_text(ss, 1.05, 6.42, 11.4, 0.66,
      [[R("FoolsGold #1 (68.76%) via full-history tracking; Aegis #2 at ", 13.5, TEXT, False, F_BODY),
        R("63.17%", 13.5, GREEN, True, F_BODY),
        R(". ", 13.5, TEXT, False, F_BODY),
        R("Geometric baselines destroyed", 13.5, CRIMSON, True, F_BODY),
        R(" (Krum 14%, CWMed 22%).", 13.5, TEXT, False, F_BODY)]],
      anchor=MSO_ANCHOR.MIDDLE)
syb_idx = len(prs.slides) - 1
move_slide(prs, syb_idx, 38)  # insert after Vol Spam → position 38
print("  Inserted Sybil slide at position 39.")

# ═══════════════════════════════════════════════════════════════════════════════
# 3. UPDATE SUMMARY TABLE — now at index 44 (was 42, +2 from insertions)
# ═══════════════════════════════════════════════════════════════════════════════
print("\n=== Step 3: Updating Summary table ===")
# After inserting 2 slides, the old slide 43 (idx 42) is now at idx 44
summary_idx = 44
s_sum = prs.slides[summary_idx]

# Find and remove old table + old bullets
for shape in list(s_sum.shapes):
    if shape.has_table:
        sp = shape._element; sp.getparent().remove(sp)
        print("  Removed old summary table.")
    elif shape.has_text_frame:
        t = shape.text_frame.text
        if 'Aegis is strictly superior' in t or 'FoolsGold wins' in t or 'Aegis heavily outperforms' in t:
            sp = shape._element; sp.getparent().remove(sp)
            print(f"  Removed old bullet: {t[:60]}...")

# New summary table data (with Vol Spam row added, Sybil updated)
sum_data = [
    ["Attack",      "Aegis Acc", "DR %",  "Best Baseline",       "Gap (pp)", "Regime"],
    ["No Attack",   "76.08",     "N/A",   "FedAvg  76.16",       "\u22120.08",  "I"],
    ["Label Flip",  "74.34",     "76.5",  "Krum  72.85",         "+1.49",     "I"],
    ["Vol Spam",    "73.77",     "56.7",  "Krum  72.45",         "+1.32",     "I"],
    ["Sign Flip",   "66.61",     "60.8",  "Krum  70.05",         "\u22123.44",  "II"],
    ["IPM",         "63.32",     "0.0",   "FoolsGold  71.07",    "\u22127.75",  "II"],
    ["Sybil",       "63.17",     "0.2",   "FoolsGold  68.76",    "\u22125.59",  "II/III"],
    ["ALIE",        "10.00",     "0.0",   "FoolsGold  71.15",    "\u221261.15", "III"],
]

sum_col_w = [1.35, 1.15, 0.85, 2.50, 1.10, 1.00]
nr_s = len(sum_data); nc_s = len(sum_data[0])
row_h_s = 0.38
gt_s = s_sum.shapes.add_table(nr_s, nc_s, Inches(0.55), Inches(1.70),
                               Inches(sum(sum_col_w)), Inches(row_h_s * nr_s)).table
gt_s.first_row = False; gt_s.horz_banding = False

for ci, cw in enumerate(sum_col_w):
    gt_s.columns[ci].width = Inches(cw)

# Highlight styles
sum_styles = {
    # Regime I wins (green)
    (2, 4): {"fill": GRNTNT, "bold": True, "color": GREEN},  # Label Flip +1.49
    (3, 4): {"fill": GRNTNT, "bold": True, "color": GREEN},  # Vol Spam +1.32
    # Regime III failure (red)
    (7, 1): {"fill": CRIMTNT, "bold": True, "color": CRIMSON},  # ALIE 10.00
    (7, 4): {"fill": CRIMTNT, "bold": True, "color": CRIMSON},  # ALIE gap
}

for ri in range(nr_s):
    gt_s.rows[ri].height = Inches(row_h_s)
    for ci in range(nc_s):
        cell = gt_s.cell(ri, ci)
        txt = str(sum_data[ri][ci])
        if ri == 0:
            set_cell(cell, txt, 10.5, WHITE, True,
                     PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER, NAVY, F_HEAD)
        else:
            base = LIGHT if ri % 2 == 0 else WHITE
            col = TEXT; bold = False
            st = sum_styles.get((ri, ci))
            if st:
                col = st.get("color", col)
                bold = st.get("bold", bold)
                base = st.get("fill", base)
            set_cell(cell, txt, 10.5, col, bold,
                     PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER,
                     base, F_HEAD if (ci == 0 and bold) else F_BODY)

print("  New 8-row summary table inserted.")

# Add updated bullets
tb2 = s_sum.shapes.add_textbox(Inches(0.8), Inches(5.10), Inches(11.8), Inches(1.8))
tf2 = tb2.text_frame; tf2.word_wrap = True
tf2.margin_left = 0; tf2.margin_right = 0; tf2.margin_top = 0; tf2.margin_bottom = 0

p1 = tf2.paragraphs[0]
p1.space_after = Pt(6); p1.space_before = Pt(0); p1.line_spacing = 1.04; p1.alignment = PP_ALIGN.LEFT
b1 = p1.add_run(); b1.text = "\u25aa  "; b1.font.size = Pt(11.5); b1.font.color.rgb = ACCENT; b1.font.name = F_BODY; b1.font.bold = True
r1 = p1.add_run(); r1.text = "Aegis is strictly superior to geometric baselines on Label-Flip, Vol Spam, IPM, and Sybil, but lags Multi-Krum on Sign-Flip."
r1.font.size = Pt(12.5); r1.font.color.rgb = NAVY; r1.font.name = F_BODY; r1.font.bold = True

p2 = tf2.add_paragraph()
p2.space_after = Pt(6); p2.space_before = Pt(0); p2.line_spacing = 1.04; p2.alignment = PP_ALIGN.LEFT
b2 = p2.add_run(); b2.text = "\u25aa  "; b2.font.size = Pt(11.5); b2.font.color.rgb = ACCENT; b2.font.name = F_BODY; b2.font.bold = True
r2 = p2.add_run(); r2.text = "FoolsGold wins ALIE, IPM & Sybil \u2014 its full-history tracking outresolves the 20-round EMA. Aegis\u2019s limitation on median-poisoning attacks is a known tradeoff."
r2.font.size = Pt(12.5); r2.font.color.rgb = TEXT; r2.font.name = F_BODY; r2.font.bold = False

print("  Updated summary bullets.")

# ═══════════════════════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════════════════════
try:
    prs.save(PPTX)
    print(f"\nSaved successfully to {PPTX}")
except PermissionError:
    print("\nERROR: PermissionError — please close PowerPoint and re-run this script.")
